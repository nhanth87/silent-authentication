#!/usr/bin/env python3
"""Build Digicom-ET Silent Auth Ethiopia-themed PowerPoint (~20 slides)."""

from pathlib import Path
import subprocess
import shutil

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
OUT = ROOT / "DigicomET_Silent_Auth_Ethiopia.pptx"

# Ethiopia palette
GREEN = RGBColor(0x07, 0x89, 0x30)
YELLOW = RGBColor(0xFC, 0xDD, 0x09)
RED = RGBColor(0xDA, 0x12, 0x1A)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
CREAM = RGBColor(0xFF, 0xF8, 0xE7)
TEAL = RGBColor(0x0D, 0x73, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GRAY = RGBColor(0x4B, 0x55, 0x63)


def set_run(run, size=18, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_flag_bar(slide, prs):
    """Top stripe: green / yellow / red."""
    w = prs.slide_width
    h = Inches(0.18)
    for i, color in enumerate([GREEN, YELLOW, RED]):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(i * w // 3), 0, Emu(w // 3 + 1000), h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()


def fill_slide_bg(slide, color=CREAM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, prs, page, total=20):
    add_textbox(
        slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
        f"Digicom-ET  ·  Silent Authentication for Ethiopian Banks  ·  {page}/{total}",
        size=10, color=GRAY, align=PP_ALIGN.LEFT,
    )


def svg_to_png(svg_path: Path, png_path: Path, scale=2):
    """Rasterize SVG → PNG via cairosvg, rsvg-convert, or inkscape."""
    if png_path.exists() and png_path.stat().st_mtime >= svg_path.stat().st_mtime:
        return png_path

    # try cairosvg
    try:
        import cairosvg
        cairosvg.svg2png(url=str(svg_path), write_to=str(png_path), scale=scale)
        return png_path
    except Exception:
        pass

    # try rsvg-convert
    if shutil.which("rsvg-convert"):
        subprocess.run(
            ["rsvg-convert", "-z", str(scale), "-o", str(png_path), str(svg_path)],
            check=True,
        )
        return png_path

    # try inkscape
    if shutil.which("inkscape"):
        subprocess.run(
            ["inkscape", str(svg_path), f"--export-filename={png_path}",
             f"--export-dpi={96 * scale}"],
            check=True,
        )
        return png_path

    # fallback: Pillow + svglib
    try:
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM
        drawing = svg2rlg(str(svg_path))
        renderPM.drawToFile(drawing, str(png_path), fmt="PNG")
        return png_path
    except Exception as e:
        raise RuntimeError(
            f"Cannot rasterize {svg_path.name}. Install cairosvg or rsvg-convert. ({e})"
        )


def add_picture(slide, name, left, top, width=None, height=None):
    svg = ASSETS / name
    png = ASSETS / (svg.stem + ".png")
    svg_to_png(svg, png)
    kwargs = {"image_file": str(png), "left": left, "top": top}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(**kwargs)


def section_title_slide(prs, title, subtitle, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_slide_bg(slide)
    add_flag_bar(slide, prs)
    # accent block
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.2), Inches(13.333), Inches(2.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1),
                title, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.7),
                subtitle, size=18, color=CREAM, align=PP_ALIGN.CENTER)
    add_footer(slide, prs, page)
    return slide


def content_slide(prs, title, bullets, page, img=None, img_w=Inches(5.5)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    add_flag_bar(slide, prs)
    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.6),
                title, size=28, bold=True, color=DARK)
    # green underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(2.5), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()

    left_w = Inches(6.5) if img else Inches(12)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), left_w, Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"•  {b}"
        set_run(run, size=16, color=DARK)

    if img:
        add_picture(slide, img, Inches(7.2), Inches(1.4), width=img_w)
    add_footer(slide, prs, page)
    return slide


def build():
    # ensure SVGs exist
    subprocess.run(
        ["python3", str(ROOT / "scripts" / "generate_svgs.py")],
        check=True,
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    titles = []

    # ---- 1 Title ----
    titles.append("Digicom-ET")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s)
    add_flag_bar(s, prs)
    panel = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), Inches(13.333), Inches(4.2)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = GREEN
    panel.line.fill.background()
    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
                "DIGICOM-ET", size=20, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2),
                "Silent Authentication for Banks", size=40, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.6),
                "Prove Chú Phỉnh's phone on the network — no SMS, under 3 seconds",
                size=18, color=CREAM, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.4),
                "Ethiopia  ·  VAS adapter on Ethio Telecom  ·  2026",
                size=14, color=YELLOW, align=PP_ALIGN.CENTER)
    add_footer(s, prs, 1)

    # ---- 2 Agenda ----
    titles.append("Agenda")
    content_slide(prs, "Agenda", [
        "Why Ethiopia — mobile banking at scale, OTP pain everywhere",
        "The SMS OTP problem & the threats banks face today",
        "Meet Chú Phỉnh — frictionless login, real security",
        "Digicom-ET Silent Auth: how the VAS adapter works",
        "Architecture on top of Ethio Telecom — not inside it",
        "Fallback SMS stays operator-billed; telco revenue protected",
        "Security, integration, roadmap, next steps",
    ], 2)

    # ---- 3 Ethiopia context ----
    titles.append("Why Ethiopia")
    content_slide(prs, "Why Ethiopia", [
        "70M+ mobile subs — banking runs on the phone",
        "SMS OTP still default: slow, expensive, easy to attack",
        "Banks need instant login without weakening fraud controls",
        "Network-side identity as a VAS layer — no core rebuild",
        "Digicom-ET: Silent Auth built for Ethiopian banks",
    ], 3, img="08_ethiopia.svg", img_w=Inches(5.2))

    # ---- 4 Problem ----
    titles.append("The SMS OTP problem")
    content_slide(prs, "The SMS OTP problem", [
        "3–30 s wait → abandoned logins and support calls",
        "SS7 / interconnect: OTP redirected before the user sees it",
        "SIM-swap: attacker owns the number, receives the code",
        "Phishing proxies harvest codes in real time",
        "Banks pay per SMS and still eat fraud losses",
    ], 4, img="01_problem_sms.svg", img_w=Inches(5.4))

    # ---- 5 Persona ----
    titles.append("Meet Chú Phỉnh")
    content_slide(prs, "Meet Chú Phỉnh", [
        "Daily mobile-banking user in Addis Ababa",
        "Wants: open app → recognized → done",
        "Refuses: wait, hunt for SMS, type six digits",
        "Must feel safe — and actually be safe — from fraud",
        "Silent Auth delivers both: speed and assurance",
    ], 5, img="02_persona.svg", img_w=Inches(5.4))

    # ---- 6 Threats ----
    titles.append("Threats Silent Auth removes")
    content_slide(prs, "Threats Silent Auth removes", [
        "SS7 intercept — no SMS payload to steal",
        "Phishing / OTP relay — nothing for the user to leak",
        "Premium-rate / AIT abuse — no OTP blast to fraud numbers",
        "Delivery failure — login not hostage to SMS latency",
        "SIM-swap — network signals flag recent handset changes",
    ], 6, img="06_threats.svg", img_w=Inches(5.3))

    # ---- 7 Solution intro ----
    titles.append("The solution")
    section_title_slide(
        prs,
        "The solution",
        "Digicom-ET Silent Auth — VAS adapter that proves the phone is on-network",
        7,
    )

    # ---- 8 How it works ----
    titles.append("How Silent Auth works")
    content_slide(prs, "How Silent Auth works", [
        "App on cellular captures IP + port + timestamp",
        "Bank backend calls Digicom-ET POST /verify",
        "Resolver binds IP:port → MSISDN via PGW session",
        "Verifier queries HLR/HSS: MAP (PSI/ATI/SAI) or Diameter S6a",
        "Policy engine: APPROVE silently or route to SMS fallback",
    ], 8, img="03_silent_flow.svg", img_w=Inches(5.6))

    # ---- 9 Before / After ----
    titles.append("Before vs After")
    content_slide(prs, "Before vs After", [
        "Before: wait → type OTP → pray it wasn't intercepted",
        "After: silent verify → session in under 3 seconds",
        "Happy path: zero SMS, zero typing",
        "Same bank app — Digicom VAS plugs in behind the API",
        "Chú Phỉnh: open app, you're in",
    ], 9, img="09_before_after.svg", img_w=Inches(5.5))

    # ---- 10 Architecture ----
    titles.append("Architecture — adapter on top")
    content_slide(prs, "Architecture — adapter on top", [
        "Digicom-ET = thin VAS layer between banks and Ethio Telecom",
        "Reuses existing PGW, HLR/HSS, SMSC — zero core rewrite",
        "Banks call HTTPS API; Digicom handles signalling orchestration",
        "Operator network remains source of truth for identity",
        "We adapt on top — we never replace the telecom stack",
    ], 10, img="04_architecture.svg", img_w=Inches(5.8))

    # ---- 11 Money / pocket ----
    titles.append("We do NOT touch Ethio Telecom's pocket")
    content_slide(prs, "We do NOT touch Ethio Telecom's pocket", [
        "Digicom sells verification to banks — not SMS wholesale",
        "A2P / interconnect SMS revenue stays 100% with the operator",
        "Signalling, data, and subscriber records remain carrier-owned",
        "Optional API revenue share — additive VAS, not cannibalization",
        "Every fallback OTP is operator-billed through Ethio Telecom SMSC",
    ], 11, img="07_money_model.svg", img_w=Inches(5.5))

    # ---- 12 Fallback ----
    titles.append("Fallback to SMS")
    content_slide(prs, "Fallback to SMS", [
        "Wi-Fi-only, timeout, or low assurance → SMS OTP safety net",
        "Fail-closed: never approve without cryptographic evidence",
        "SMS is backup — silent path is primary",
        "Digicom orchestrates policy only; SMS rides operator rails",
        "Fallback traffic = operator revenue, not Digicom bypass",
    ], 12, img="05_fallback.svg", img_w=Inches(5.2))

    # ---- 13 Defense in depth ----
    titles.append("Defense in depth")
    content_slide(prs, "Defense in depth", [
        "Layer 1: Silent Auth eliminates most OTP exposure",
        "Layer 2: SMS Home Routing / firewall for residual OTP",
        "Layer 3: Operator SS7, Diameter, 5G SEPP border controls",
        "Digicom orchestrates bank-facing identity policy",
        "Aligned with GSMA FS.11 / FS.19 / FS.21 / FS.36",
    ], 13, img="10_shield.svg", img_w=Inches(5.0))

    # ---- 14 Value for banks ----
    titles.append("Value for Ethiopian banks")
    content_slide(prs, "Value for Ethiopian banks", [
        "Higher conversion — fewer abandoned logins",
        "Lower OTP SMS spend on the silent happy path",
        "Stronger fraud posture vs SS7, phishing, SIM-swap",
        "Drop-in integration: HTTPS API + mobile SDK",
        "Risk-based step-up for high-value transactions",
    ], 14)

    # ---- 15 Value for telecom ----
    titles.append("Value for Ethio Telecom partnership")
    content_slide(prs, "Value for Ethio Telecom partnership", [
        "New enterprise VAS revenue — SMS wholesale untouched",
        "Banks stay on cellular; operator keeps subscriber engagement",
        "Digicom sells to banks; telco supplies network truth",
        "Shared policy + optional per-verify API revenue share",
        "Positions Ethiopia for Open Gateway / CAMARA leadership",
    ], 15)

    # ---- 16 Security highlights ----
    titles.append("Security highlights")
    content_slide(prs, "Security highlights", [
        "Intra-network HLR/HSS queries only (FS.11 Cat. 1 on interconnect)",
        "IP + port + timestamp — CGNAT-safe binding",
        "SIM-swap freshness via SAI / lastUpdateLocation",
        "Hard timeouts: resolver 300 ms, MAP/Diameter 2 s, total ≤ 3 s",
        "MSISDN never exposed to the app — bank backend only",
    ], 16)

    # ---- 17 Integration ----
    titles.append("Bank integration")
    content_slide(prs, "Bank integration", [
        "1. Embed Digicom SDK in the mobile banking app",
        "2. On login: app sends deviceCred + cellular binding to bank",
        "3. Bank backend calls Digicom POST /verify",
        "4. Digicom returns match + assurance level (or FALLBACK)",
        "5. Bank issues session — or triggers operator-billed SMS OTP",
    ], 17)

    # ---- 18 Roadmap ----
    titles.append("Roadmap")
    content_slide(prs, "Roadmap", [
        "Phase 1: Pilot 1–2 banks — IP-match + MAP verifier",
        "Phase 2: Diameter S6a + SIM-swap signal API",
        "Phase 3: TS.43 / Wi-Fi silent path for wider coverage",
        "Phase 4: CAMARA Number Verification alignment",
        "Ongoing: SMS Home Routing for operator-protected fallback",
    ], 18)

    # ---- 19 CTA ----
    titles.append("Let's make Chú Phỉnh's login silent")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s)
    add_flag_bar(s, prs)
    panel = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.5)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = TEAL
    panel.line.fill.background()
    add_textbox(s, Inches(1.8), Inches(2.2), Inches(9.7), Inches(0.8),
                "Let's make Chú Phỉnh's login silent", size=32, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.8), Inches(3.2), Inches(9.7), Inches(1.2),
                "Pilot Digicom-ET Silent Auth with your bank.\n"
                "VAS adapter on Ethio Telecom — telco SMS revenue stays intact.",
                size=16, color=CREAM, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.8), Inches(4.5), Inches(9.7), Inches(0.5),
                "contact@digicom-et.example  ·  Addis Ababa",
                size=14, color=YELLOW, align=PP_ALIGN.CENTER)
    add_footer(s, prs, 19)

    # ---- 20 Thank you ----
    titles.append("Ameseginalehu · Thank you")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s)
    add_flag_bar(s, prs)
    add_textbox(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1),
                "Ameseginalehu  ·  Thank you", size=40, bold=True,
                color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.6),
                "Digicom-ET — VAS Silent Authentication for Ethiopian Banks",
                size=18, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5),
                "Adapter on top · Operator revenue protected · Built for Ethiopia",
                size=14, color=TEAL, align=PP_ALIGN.CENTER)
    add_footer(s, prs, 20)

    assert len(titles) == 20, f"Expected 20 slides, got {len(titles)}"
    prs.save(OUT)
    print(f"Saved → {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Titles:")
    for i, t in enumerate(titles, 1):
        print(f"  {i:2d}. {t}")


if __name__ == "__main__":
    build()
