#!/usr/bin/env python3
"""Digicom-ET Silent Auth — Sea theme proposal deck ("nước xanh và trắng").

21 slides, dominated by deep-sea blues and white. Call flows drawn natively with
pptx shapes matching docs/design/ (silent-auth-flow.md, camara-flow-diagram.svg,
unified-identity-sms-security-architecture.md) plus a micro-jainslee (JSR-240 /
JAIN SLEE 1.1, a.k.a SR-240) tech-foundation intro.

Usage:  python3 slides/scripts/build_pptx_sea.py
Output: slides/DigicomET_Silent_Auth_Proposal_Sea.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "DigicomET_Silent_Auth_Proposal_Sea.pptx"
TOTAL = 21

# ---- Sea palette: "nước xanh và trắng" (deep water blues + white) ----
NAVY = RGBColor(0x04, 0x32, 0x55)      # deep sea
OCEAN = RGBColor(0x0E, 0x6B, 0xA8)     # ocean blue
AQUA = RGBColor(0x2F, 0xA4, 0xC6)      # light teal / sea surface
CYAN = RGBColor(0x7F, 0xD3, 0xEC)      # foam highlight
WAVE = RGBColor(0xBF, 0xE4, 0xF0)      # pale wave
FOAM = RGBColor(0xEA, 0xF6, 0xFB)      # near-white sea foam
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x14, 0x2B, 0x3C)      # navy-ink text
GRAY = RGBColor(0x4D, 0x5D, 0x6A)
SAND = RGBColor(0xE8, 0xC8, 0x7A)      # tiny warm accent (sparingly)


def set_run(run, size=16, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=16, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def fill_bg(slide, color=FOAM):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_wave_bottom(slide, color=WAVE, top=Inches(6.9)):
    """Three overlapping ovals at the bottom edge suggest sea water."""
    for cx, w in [(1.0, 2.6), (4.2, 3.4), (8.3, 3.0), (11.9, 2.3)]:
        shp = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - w / 2), top, Inches(w), Inches(1.1)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()


def add_footer(slide, prs, page):
    add_textbox(
        slide, Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.28),
        f"Digicom-ET  ·  Silent Authentication Proposal (Sea)  ·  {page}/{TOTAL}",
        size=10, color=GRAY,
    )
    add_wave_bottom(slide, WAVE, top=Inches(7.28))


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.55),
                title, size=26, bold=True, color=NAVY)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(0.85), Inches(2.4), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = AQUA
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.93), Inches(12.3), Inches(0.3),
                    subtitle, size=12, color=GRAY)


def bullets(slide, left, top, width, height, items, size=14, gap=8,
            color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size=size, color=color)


def add_box(slide, left, top, width, height, text, fill=FOAM, line=OCEAN,
            t_color=DARK, size=12, bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
            line_w=1.0, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=t_color)
    return shp


def add_multi_box(slide, left, top, width, height, lines, fill=FOAM, line=OCEAN,
                  t_color=DARK, size=11, bold=False, line_w=1.0):
    """Rounded box with several small centered lines."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = txt
        set_run(run, size=size, bold=bold and i == 0, color=t_color)
    return shp


def set_line(line, color, width=2.0, dash=False):
    line.color.rgb = color
    line.width = Pt(width)
    ln = line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(d)


def add_arrow(slide, x1, y1, x2, y2, color=OCEAN, width=2.0, dash=False,
              head=True):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2
    )
    set_line(conn.line, color, width=width, dash=dash)
    if head:
        ln = conn.line._get_or_add_ln()
        tail = ln.makeelement(
            qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}
        )
        ln.append(tail)
    return conn


def add_circle_num(slide, cx, cy, num, fill=OCEAN, r=0.16):
    c = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Emu(int(cx - Inches(r))), Emu(int(cy - Inches(r))),
        Inches(2 * r), Inches(2 * r)
    )
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    tf = c.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(num)
    set_run(run, size=11, bold=True, color=WHITE)
    return c


def add_table_sea(slide, headers, rows, left, top, width, col_widths, row_h=0.34):
    """Sea-styled table: navy head, alternating foam/white rows."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, Inches(row_h * n_rows)
    ).table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        set_run(run, size=11, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        bg = FOAM if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            set_run(run, size=10, color=DARK)
    return table


def _label_arrow(slide, x1, y1, x2, y2, label, size=10, dy=-0.22, color=DARK):
    """Small centered label above a horizontal arrow (y1==y2)."""
    cx1, cx2 = min(x1, x2), max(x1, x2)
    w = max(cx2 - cx1, Inches(1.1))
    mid = Emu(int((x1 + x2) / 2))
    add_textbox(slide, Emu(int(mid - w / 2)), Emu(int(y1 + dy)),
                w, Inches(0.3), label, size=size, color=color,
                align=PP_ALIGN.CENTER)


def seq_diagram(slide, left, top, height, participants, steps):
    """Minimal native sequence diagram.

    participants: list of (title, subtitle)
    steps: list of (from_idx, to_idx, label, dashed)
    All positions are Emu ints (left/top/height from Inches()). Lifelines run
    from top to top+height; arrow rows spaced evenly down the diagram.
    """
    n = len(participants)
    n_steps = len(steps)
    span = Inches(11.2)
    xs = [Emu(int(left + span * i // (n - 1))) for i in range(n)]
    head_h = Inches(1.0)
    row_h = Emu(int((height - head_h) / max(n_steps, 1)))

    # participant headers
    for i, (title, sub) in enumerate(participants):
        w = Inches(1.9)
        hx = Emu(int(xs[i] - w / 2))
        add_box(slide, hx, top, w, Inches(0.62), title, fill=NAVY,
                t_color=WHITE, size=11, bold=True, line=NAVY)
        if sub:
            add_textbox(slide, Emu(int(xs[i] - w / 2)), Emu(int(top + 0.6)),
                        w, Inches(0.3), sub, size=8, color=GRAY,
                        align=PP_ALIGN.CENTER)
        add_arrow(slide, xs[i], Emu(int(top + 0.62)), xs[i],
                  Emu(int(top + height)), color=CYAN, width=1.2, head=False)
    # steps
    y = Emu(int(top + head_h))
    for (a, b, label, dashed) in steps:
        ya = Emu(int(y + head_h / 3))
        if a == b:
            # self message: small loop on the right
            add_arrow(slide, xs[a], ya, Emu(int(xs[a] + Inches(0.28))), ya,
                      color=AQUA, width=1.6, dash=dashed)
            add_textbox(slide, Emu(int(xs[a] + Inches(0.32))),
                        Emu(int(ya - Inches(0.12))), Inches(2.2), Inches(0.3),
                        label, size=9, color=DARK)
        else:
            add_arrow(slide, xs[a], ya, xs[b], ya, color=OCEAN, width=1.8,
                      dash=dashed)
            _label_arrow(slide, xs[a], ya, xs[b], ya, label, size=9)
        y = Emu(int(y + row_h))
    return xs


def slide_cover(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, NAVY)
    # layered waves at the bottom of the deep sea
    add_wave_bottom(s, OCEAN, top=Inches(5.6))
    add_wave_bottom(s, AQUA, top=Inches(6.0))
    add_wave_bottom(s, CYAN, top=Inches(6.4))
    add_wave_bottom(s, WHITE, top=Inches(6.85))
    add_textbox(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4),
                "DIGICOM-ET  ·  PROPOSAL  ·  ETHIOPIA", size=16, bold=True,
                color=CYAN, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.5),
                "Silent Authentication", size=54, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(3.15), Inches(11.7), Inches(0.5),
                "SMS OTP lặng lẽ biến mất — phone-on-network là mật khẩu",
                size=20, color=CYAN, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.9),
                "CAMARA Number Verification  ·  TS.43 EAP-AKA  ·  MAP/Diameter verifier\n"
                "Chạy trên nền micro-jainslee (SR-240 / JAIN SLEE 1.1)",
                size=13, color=WAVE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
                "confidential  ·  designed & engineered by nhanth87 (Tran Nhan)",
                size=11, color=WAVE, align=PP_ALIGN.CENTER)
    return s


def slide_microjainslee(prs, page):
    """Tech foundation: micro-jainslee, the SR-240 / JSR-240 / JAIN SLEE 1.1 engine."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Nền tảng công nghệ — micro-jainslee (SR-240)",
              "JAIN SLEE 1.1 / JSR-240 — container SLEE nhúng: event, activity, timer, RA, SBB")
    add_textbox(s, Inches(0.5), Inches(1.25), Inches(6.4), Inches(0.35),
                "micro-jainslee là gì", size=15, bold=True, color=OCEAN)
    bullets(s, Inches(0.5), Inches(1.6), Inches(6.45), Inches(4.5), [
        "Container SLEE R&D (~3 KLOC core) nhúng trong Quarkus / Spring / Jakarta — không cần server App đặt riêng",
        "Triển khai subset JAIN SLEE 1.1 (JSR-240): SBB lifecycle, ActivityContext, EventRouting, Timer Facility (§9), RA SPI (§13)",
        "EventRouter chạy LMAX Disruptor — nhận event, route đúng activity, gọi SBB entity theo vòng đời",
        "Timer: Netty HashedWheelTimer qua SleeTimerSchedulerBridge — giới hạn dialog & timeout fail-closed",
        "Java 25: virtual-thread entity pool + hướng ZGC — sinh ra để chạy trong JVM nhỏ của SAS",
        "APT sinh GeneratedEventTypes + SbbIndexLoader — map event → SBB tự động khi start",
    ], size=13, gap=9)
    add_box(s, Inches(0.5), Inches(6.15), Inches(6.45), Inches(0.75),
            "Chỉ dành cho R&D & nhúng sản phẩm kiểu SAS lab — "
            "production USSD 7.3 vẫn dùng Mobicents SLEE",
            fill=RGBColor(0xFB, 0xF3, 0xE0), line=SAND, t_color=DARK, size=12)
    add_textbox(s, Inches(7.1), Inches(1.25), Inches(5.7), Inches(0.35),
                "Kiến trúc module & luồng event", size=15, bold=True, color=OCEAN)
    add_multi_box(s, Inches(7.1), Inches(1.7), Inches(2.75), Inches(0.85),
                  ["App (Quarkus)", "REST / services"], fill=WAVE, line=OCEAN,
                  size=10, bold=True)
    add_arrow(s, Inches(9.85), Inches(2.12), Inches(10.25), Inches(2.12))
    add_multi_box(s, Inches(10.25), Inches(1.7), Inches(2.6), Inches(0.85),
                  ["micro-jainslee", "MicroSleeConfiguration"], fill=NAVY,
                  line=NAVY, t_color=WHITE, size=10, bold=True)
    y = 2.75
    for head, sub in [
        ("jainslee-api", "SbbContext · Event · RA SPI"),
        ("jainslee-core", "ActivityContext · Disruptor EventRouter"),
        ("jainslee-scheduler", "HashedWheelTimer → SleeTimerBridge"),
        ("jainslee-apt", "GeneratedEventTypes · SbbIndexLoader"),
    ]:
        add_multi_box(s, Inches(10.25), Inches(y), Inches(2.6), Inches(0.72),
                      [head, sub], fill=FOAM, line=AQUA, size=10)
        y += 0.92
    add_multi_box(s, Inches(7.1), Inches(2.75), Inches(2.6), Inches(3.2),
                  ["SBB (VerifySbb)", "onEvent → FSM\nfireEvent ra tiếp",
                   "CMP · ChildSBB", "RA wrapper: endpoint + command",
                   "bootstrap.fireEvent()"], fill=WAVE, line=OCEAN,
                  size=10)
    add_textbox(s, Inches(7.1), Inches(6.15), Inches(5.75), Inches(0.75),
                "Quy tắc: mọi I/O signalling nằm trong RA delegate (seam /ras/); "
                "không có executor/timer/socket tay ngoài container (gate H24).",
                size=11, color=GRAY)
    add_footer(s, prs, page)
    return s


def slide_h24(prs, page):
    """H24: the SAS is one thin northbound over a single micro-jainslee runtime."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "H24 — chỉ micro-jainslee chạy SAS, không code vòng quanh container",
              "HTTP layer chỉ nộp event và chờ kết quả; trạng thái / timer / signalling đều trong container")
    add_multi_box(s, Inches(0.4), Inches(2.4), Inches(1.9), Inches(1.5),
                  ["Bank Backend", "POST /verify\n(mTLS, reqId)"], fill=WAVE,
                  line=OCEAN, size=10, bold=True)
    add_arrow(s, Inches(2.3), Inches(3.0), Inches(2.85), Inches(3.0))
    add_multi_box(s, Inches(2.85), Inches(2.4), Inches(2.35), Inches(1.5),
                  ["HTTP seam (Quarkus)", "VerifyResource\nchỉ submit event\nchờ CompletableFuture"],
                  fill=WAVE, line=OCEAN, size=10, bold=True)
    add_arrow(s, Inches(5.2), Inches(3.0), Inches(5.75), Inches(3.0))
    big = Emu(int(Inches(5.75)))
    big_w = Inches(4.7)
    add_multi_box(s, big, Inches(1.6), big_w, Inches(4.7),
                  ["micro-jainslee container (1 container)", ""],
                  fill=NAVY, line=NAVY, t_color=WHITE, size=12, bold=True)
    inner_left = Emu(int(big + Inches(0.35)))
    inner_w = Inches(4.0)
    add_multi_box(s, inner_left, Inches(2.2), inner_w, Inches(1.5),
                  ["RA delegates — transport seam (/ras/)",
                   "resolver-ra · map-verifier-ra\ns6a-verifier-ra · swx-verifier-ra\n"
                   "chỉ nơi này mới có I/O SS7/Diameter — bootstrap.fireEvent()"],
                  fill=FOAM, line=AQUA, t_color=DARK, size=10)
    add_multi_box(s, inner_left, Inches(4.0), inner_w, Inches(1.7),
                  ["VerifySbb — onEvent → FSM",
                   "RESOLVING → VERIFYING → SCORING\n→ APPROVED · FALLBACK (fail-closed)\n"
                   "timer dialog có giới hạn, timeout ⇒ abort"],
                  fill=FOAM, line=AQUA, t_color=DARK, size=10)
    add_arrow(s, Emu(int(big + Inches(2.2))), Inches(3.7),
              Emu(int(big + Inches(2.2))), Inches(4.0))
    add_multi_box(s, Inches(11.0), Inches(2.4), Inches(1.9), Inches(1.5),
                  ["SAS outcome", "APPROVED | FALLBACK\nassurance + reqId"], fill=WAVE,
                  line=OCEAN, size=10, bold=True)
    add_textbox(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.75),
                "Ràng buộc kiến trúc (gate H24): com.microjainslee.core.* chỉ truy cập từ SasBootstrap; "
                "cấm javax.slee/jakarta.slee (không container thứ hai); cấm executor/timer/socket tay ngoài /ras/; "
                "RA không được bypass container. — Kiểm tra tự động qua harness.",
                size=11, color=GRAY)
    add_footer(s, prs, page)
    return s


def slide_agenda(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Agenda — đề xuất Silent Authentication cho ngân hàng Ethiopia")
    items = [
        ("01", "Vì sao ngân hàng cần Silent Auth ngay bây giờ", "SMS OTP tấn công & chi phí banking"),
        ("02", "Giải pháp: Xác thực im lặng hai tầng", "Resolver + Verifier, IP-match & TS.43"),
        ("03", "Hạ tầng công nghệ", "micro-jainslee (SR-240) chạy SAS — H24 boundary"),
        ("04", "CAMARA Number Verification", "Giao diện /verify cho bank, token + mTLS"),
        ("05", "Assurance & an toàn", "Điểm tin cậy, fail-closed, không ATI liên mạng"),
        ("06", "Kiến trúc thống nhất A+B", "Silent Auth + bảo vệ kênh SMS/SS7"),
        ("07", "Mô hình thương mại & lộ trình", "Digicom-ET bán VAS, triển khai 4 pha"),
    ]
    y = 1.35
    for num, title, sub in items:
        add_circle_num(s, Inches(0.85), Inches(y + 0.28), int(num), fill=OCEAN)
        add_textbox(s, Inches(1.3), Inches(y), Inches(8.4), Inches(0.4),
                    title, size=16, bold=True, color=NAVY)
        add_textbox(s, Inches(1.3), Inches(y + 0.38), Inches(11.4), Inches(0.3),
                    sub, size=11, color=GRAY)
        y += 0.8
    add_footer(s, prs, page)
    return s


def slide_why(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Vì sao Ethiopia, vì sao ngay bây giờ",
              "Mobile money phổ cập nhanh — nhưng lớp xác thực vẫn là SMS OTP dễ đánh cắp")
    bullets(s, Inches(0.5), Inches(1.35), Inches(6.2), Inches(5.2), [
        "Mobile-first: phần lớn người dùng tài chính vào ngân hàng qua app trên 2G/3G/4G — không bàn phím phần cứng, không FIDO keys",
        "SMS OTP ngày càng rẻ để tấn công: SS7 interception, SIM swap, malware trên thiết bị",
        "Ngân hàng chịu 100% chi phí cho luồng OTP: SMSC, trunk, hỗ trợ khách hàng khi gửi trễ/thất bại",
        "Chính phủ & NBE đang siết yêu cầu an toàn: bằng chứng thiết bị + chống gian lận chuyển khoản",
        "Digicom-ET nằm sẵn trên hạ tầng operator (Ethio Telecom) — không phải nhập hạ tầng mới",
    ], size=14, gap=10)
    # right stat card
    add_multi_box(s, Inches(7.1), Inches(1.5), Inches(5.6), Inches(2.3),
                  ["Bối cảnh 2026", "Phát hành e-wallet / tài khoản số tăng theo chính sách tài chính toàn diện",
                   "Ngân hàng cạnh tranh bằng UX: login không mật khẩu rõ ràng, siêu mượt"],
                  fill=WAVE, line=OCEAN, size=12, bold=True)
    add_multi_box(s, Inches(7.1), Inches(4.1), Inches(5.6), Inches(2.3),
                  ["Vấn đề của OTP", "• Chiếm đoạt session qua SMS intercept / SIM swap",
                   "• Chi phí vận hành trọn đời rất lớn", "• Friction: gõ 6 số, gửi lại, hết hạn"],
                  fill=FOAM, line=AQUA, size=12, bold=True)
    add_footer(s, prs, page)
    return s


def slide_otp_problem(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "SMS OTP — khâu yếu nhất của ngân hàng Ethiopia",
              "Tấn công không cần hacker vào app — chỉ cần lấy được mã một lần")
    rows = [
        ("SS7 / Roaming", "Đọc mã qua quá giang liên mạng — intercept SMS trên đường đến thiết bị", "Cat 1 · FS.11"),
        ("SIM swap", "Chuyển số → SIM mới; OTP gửi cho kẻ gian, khóa tài khoản nạn nhân", "FS.11 SIM-swap"),
        ("Malware / phishing", "RAT trên điện thoại đọc SMS; OTP bị chặn trước khi người dùng nhập", "OWASP"),
        ("SMS flooding / trễ", "OTP không đến — khách hàng bỏ dở, hỗ trợ tốn kém, doanh số mất", "UX / chi phí"),
        ("Trunk/subscription", "Nhà cung cấp SMS bán số OTP; gửi vào vùng không có bảo vệ", "FS.21"),
    ]
    add_table_sea(s, ["Tấn công", "Cơ chế", "Phân loại"], rows, Inches(0.5), Inches(1.3),
                  Inches(12.3), [Inches(2.0), Inches(8.3), Inches(2.0)], row_h=0.9)
    add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                "→ Silent Auth loại bỏ mã một lần khỏi phương trình: chứng minh chính số điện thoại / SIM đang nắm trong tay.",
                size=13, bold=True, color=OCEAN)
    add_footer(s, prs, page)
    return s


def slide_solution(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Giải pháp — Silent Authentication (xác thực im lặng)",
              "Chứng minh chính số máy đang nắm trong tay — không nhập OTP, không nhập lại mật khẩu")
    # two methods side by side
    add_multi_box(s, Inches(0.5), Inches(1.45), Inches(6.0), Inches(3.3),
                  ["Phương pháp IP-match (bearer)", "Root of trust: IP↔MSISDN qua PGW + MAP/Diameter",
                   "Cần cellular data (2G/3G/4G/5G)", "CGNAT → phải có IP + port + timestamp",
                   "Fails khi Wi-Fi / không binding / stale"], fill=WAVE,
                  line=OCEAN, size=12, bold=True)
    add_multi_box(s, Inches(6.85), Inches(1.45), Inches(6.0), Inches(3.3),
                  ["Phương pháp SIM / TS.43 (EAP-AKA)", "Root of trust: chính credential SIM",
                   "Không cần cellular data — Wi-Fi + browser OK", "Shrinks fallback-to-OTP surface",
                   "Entitlement server + SWm/SWx (TS 29.273, 33.402)"], fill=FOAM,
                  line=AQUA, size=12, bold=True)
    add_box(s, Inches(0.5), Inches(5.0), Inches(12.35), Inches(1.0),
            "App-facing: CAMARA Number Verification (NV2) — bank gọi POST /verify, "
            "SAS trả {match, assurance}.", fill=NAVY, t_color=WHITE, size=13,
            bold=True)
    add_box(s, Inches(0.5), Inches(6.1), Inches(12.35), Inches(0.7),
            "Cả hai đều fail-closed: thiếu bằng chứng → FALLBACK (Passkey / TOTP / firewalled SMS OTP) — không bao giờ soft-pass.",
            fill=RGBColor(0xFB, 0xF3, 0xE0), line=SAND, t_color=DARK, size=12)
    add_footer(s, prs, page)
    return s


def slide_twostage(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Thiết kế nền: hai tầng bắt buộc — Resolver rồi Verifier",
              "MAP / Diameter không bao giờ ánh xạ IP → MSISDN; tách trách nhiệm để hết nguy cơ sai nguồn")
    # stage boxes left-to-right
    add_multi_box(s, Inches(0.5), Inches(1.9), Inches(2.6), Inches(2.6),
                  ["INPUT", "srcIP : srcPort : ts\n(+ claimedMSISDN?)", "point-in-time from device"],
                  fill=WAVE, line=OCEAN, size=12, bold=True)
    add_arrow(s, Inches(3.1), Inches(3.1), Inches(3.85), Inches(3.1))
    add_multi_box(s, Inches(3.85), Inches(1.6), Inches(2.6), Inches(3.0),
                  ["RESOLVER", "PGW / GGSN / PCRF / CGNAT", "IP+port+ts → MSISDN / IMSI",
                   "bao nhiêu MSISDN trên IP đó?"], fill=FOAM, line=AQUA,
                  size=12, bold=True)
    add_arrow(s, Inches(6.45), Inches(3.1), Inches(7.2), Inches(3.1))
    add_multi_box(s, Inches(7.2), Inches(1.6), Inches(2.6), Inches(3.0),
                  ["VERIFIER", "MAP (PSI/ATI/SAI) · Diam S6a", "MSISDN/IMSI có live? SIM-swap?",
                   "location plausibility"], fill=FOAM, line=AQUA, size=12,
                  bold=True)
    add_arrow(s, Inches(9.8), Inches(3.1), Inches(10.55), Inches(3.1))
    add_multi_box(s, Inches(10.55), Inches(1.9), Inches(2.3), Inches(2.6),
                  ["ASSURANCE", "Policy score", "APPROVED nếu đủ,\nFALLBACK nếu thiếu"],
                  fill=NAVY, line=NAVY, t_color=WHITE, size=12, bold=True)
    add_textbox(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.7),
                "Quy tắc bất biến (design invariant)\n"
                "• MAP/Diameter không ánh xạ IP→MSISDN — chỉ Resolver (PGW/PCRF/CGNAT) làm việc đó.\n"
                "• Verifier chỉ hỏi chính HLR/HSS nhà mình — không ATI liên mạng (FS.11 Cat 1).\n"
                "• Mỗi giai đoạn đúng một dialog; timeout có giới hạn ⇒ abort, không để rò dialog.",
                size=13, color=DARK)
    add_footer(s, prs, page)
    return s


def slide_e2e(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Luồng gọi E2E — Silent Auth thay OTP khi login ngân hàng",
              "IP-match path (cellular); tổng ngân sách SAS 3 giây; bất kỳ thiếu bằng chứng nào → FALLBACK")
    participants = [
        ("User", "người dùng"),
        ("Bank App", "app trên di động"),
        ("Bank Backend", "khách hàng"),
        ("SAS (micro-jainslee)", "Resolver+Verifier+Policy"),
        ("IP Resolver", "PGW/PCRF/CGNAT"),
        ("Verifier", "MAP/Diam verifier"),
        ("HLR/HSS", "chính operator"),
    ]
    steps = [
        (0, 1, "mở app · chạm login", False),
        (1, 1, "collect {srcIP, srcPort, ts, claimedMSISDN?}", False),
        (1, 2, "POST /login (deviceCred)", False),
        (2, 3, "POST /verify — mTLS + reqId", False),
        (3, 4, "resolve(srcIP, srcPort, ts)", False),
        (4, 3, "{msisdn, imsi, bearerAge}", True),
        (3, 5, "verify(msisdn, imsi)", False),
        (5, 6, "MAP PSI/ATI · Diam S6a AIR/IDR", False),
        (6, 5, "subscriberState, VLR/MME, lastUpdate", True),
        (5, 3, "{reachable, notSimSwapped, locPlausible}", True),
        (3, 3, "Policy score ≥ threshold?", False),
        (3, 2, "{match:true, assurance:HIGH, reqId}", True),
        (2, 1, "Login OK — không cần OTP", True),
        (1, 0, "phiên đã xác thực", True),
    ]
    seq_diagram(s, Inches(0.95), Inches(1.35), Inches(5.1), participants, steps)
    add_textbox(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
                "Dialog-anchor: Resolver ≤ 300ms · MAP PSI/ATI ≤ 2s · Diameter S6a ≤ 2s · tổng SAS ≤ 3s. "
                "Timeout ⇒ abort dialog → FALLBACK.  MSISDN/IMSI không bao giờ trả về app.",
                size=11, color=GRAY)
    add_footer(s, prs, page)
    return s


def slide_camara(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "CAMARA Number Verification — vùng trách nhiệm 3 phía",
              "CV: CAMARA v2.1.0 /verify — token JWS ≤ 300s, single-use; CDR ghi full flow")
    # zones
    add_box(s, Inches(0.3), Inches(1.15), Inches(3.0), Inches(4.9),
            "", fill=WAVE, line=OCEAN, shape=MSO_SHAPE.RECTANGLE, line_w=0)
    add_textbox(s, Inches(0.35), Inches(1.2), Inches(2.9), Inches(0.3),
                "BANK SIDE", size=12, bold=True, color=OCEAN)
    add_multi_box(s, Inches(0.5), Inches(1.6), Inches(2.6), Inches(1.2),
                  ["Mobile App (UE)", "SIM + access-tech\niOS/Android/Web SDK"], fill=WHITE,
                  line=OCEAN, size=10)
    add_multi_box(s, Inches(0.5), Inches(3.2), Inches(2.6), Inches(1.2),
                  ["Bank Backend", "POST /login → /verify\nmTLS · API key · reqId"], fill=WHITE,
                  line=OCEAN, size=10)
    add_box(s, Inches(3.45), Inches(1.15), Inches(6.45), Inches(4.9),
            "", fill=FOAM, line=AQUA, shape=MSO_SHAPE.RECTANGLE, line_w=0)
    add_textbox(s, Inches(3.5), Inches(1.2), Inches(6.3), Inches(0.3),
                "DIGICOM-ET — CAMARA GATEWAY ★ SẢN PHẨM BÁN CHO BANK", size=12,
                bold=True, color=OCEAN)
    add_multi_box(s, Inches(3.6), Inches(1.6), Inches(3.0), Inches(1.1),
                  ["Auth Server", "/bc-authorize · /token\nJWS ≤ 300s single-use"], fill=WHITE,
                  line=OCEAN, size=10)
    add_multi_box(s, Inches(6.8), Inches(1.6), Inches(3.0), Inches(1.1),
                  ["NV API · CAMARA v2.1.0", "POST /verify\nso claimed vs bound"], fill=WHITE,
                  line=OCEAN, size=10)
    add_multi_box(s, Inches(3.6), Inches(3.2), Inches(6.2), Inches(2.5),
                  ["Network Adapter Layer (micro-jainslee — SR-240)",
                   "Resolver: PCRF Gx/RADIUS/CGNAT · Verifier: MAP(SIP) / Diam S6a / SWx TS.43",
                   "FSM fail-closed · assurance score · CDR · dialog timeout"],
                  fill=NAVY, line=NAVY, t_color=WHITE, size=11, bold=True)
    add_box(s, Inches(10.05), Inches(1.15), Inches(3.0), Inches(4.9),
            "", fill=WAVE, line=OCEAN, shape=MSO_SHAPE.RECTANGLE, line_w=0)
    add_textbox(s, Inches(10.1), Inches(1.2), Inches(2.9), Inches(0.3),
                "ETHIO TELECOM CORE", size=12, bold=True, color=OCEAN)
    add_multi_box(s, Inches(10.25), Inches(1.6), Inches(2.6), Inches(1.1),
                  ["HLR / HSS / UDM", "MAP · S6a · SWx"], fill=WHITE, line=OCEAN, size=10)
    add_multi_box(s, Inches(10.25), Inches(3.0), Inches(2.6), Inches(1.1),
                  ["PCRF / AAA", "Gx · RADIUS"], fill=WHITE, line=OCEAN, size=10)
    add_multi_box(s, Inches(10.25), Inches(4.4), Inches(2.6), Inches(1.1),
                  ["PGW / CGNAT", "RADIUS accounting · log"], fill=WHITE, line=OCEAN, size=10)
    # flows
    add_arrow(s, Inches(3.1), Inches(2.1), Inches(3.55), Inches(2.1))
    add_arrow(s, Inches(3.1), Inches(3.8), Inches(3.55), Inches(3.8))
    add_arrow(s, Inches(6.6), Inches(2.15), Inches(6.75), Inches(2.15))
    add_arrow(s, Inches(6.7), Inches(3.3), Inches(9.95), Inches(3.3), dash=True)
    add_arrow(s, Inches(6.7), Inches(3.6), Inches(9.95), Inches(3.6), dash=True)
    add_arrow(s, Inches(6.7), Inches(3.9), Inches(9.95), Inches(3.9), dash=True)
    add_textbox(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8),
                "Giá trị: bank thấy HTTPS chuẩn Open Gateway — toàn bộ SS7/Diameter nằm sau "
                "lớp adapter của Digicom; HSS/HLR là tài nguyên operator. Không lộ MSISDN/IMSI cho app.",
                size=12, color=DARK)
    add_footer(s, prs, page)
    return s


def slide_fsm(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "SAS FSM — fail-closed, không bao giờ soft-pass",
              "Mọi ngõ thiếu bằng chứng / timeout / SIM-swap nghi ngờ / điểm thấp ⇒ FALLBACK")
    states = [
        ("RESOLVING", "IP binding + ts"),
        ("VERIFYING", "MAP/Diam/SWx"),
        ("SCORING", "weights + threshold"),
        ("APPROVED", "assurance trả về bank"),
    ]
    xs = [0.5, 3.7, 6.9, 10.1]
    for x, (nm, sub) in zip(xs, states):
        fill = NAVY if nm == "APPROVED" else FOAM
        add_multi_box(s, Inches(x), Inches(1.7), Inches(2.7), Inches(1.3),
                      [nm, sub], fill=fill, line=OCEAN if nm != "APPROVED" else NAVY,
                      size=14, bold=True)
        if x < 10.1:
            add_arrow(s, Inches(x + 2.7), Inches(2.35), Inches(x + 3.7), Inches(2.35))
    add_multi_box(s, Inches(5.4), Inches(3.6), Inches(2.6), Inches(1.1),
                  ["FALLBACK", "Passkey · TOTP · firewalled SMS OTP"], fill=RGBColor(0xFB, 0xF3, 0xE0),
                  line=SAND, size=12, bold=True, t_color=DARK)
    for x in [xs[0], xs[1], xs[2]]:
        add_arrow(s, Inches(x + 1.3), Inches(3.0), Inches(6.3), Inches(3.6),
                  color=SAND, width=1.6, dash=True)
    rows = [
        ("Resolver", "≤ 300 ms", "không binding / >1 MSISDN / stale → FALLBACK"),
        ("MAP PSI/ATI (2G/3G)", "≤ 2 s", "abort dialog, FALLBACK"),
        ("Diameter S6a (4G/5G)", "≤ 2 s", "abort dialog, FALLBACK"),
        ("Tổng SAS", "≤ 3 s", "bank dùng ngân sách login bình thường"),
    ]
    add_table_sea(s, ["Giai đoạn", "Ngân sách", "Khi hết hạn"], rows, Inches(0.5), Inches(5.0),
                  Inches(12.3), [Inches(3.2), Inches(1.6), Inches(7.5)], row_h=0.42)
    add_footer(s, prs, page)
    return s


def slide_code_flow(prs, page):
    """SAS code flow: HTTP → event → container → RA command → response."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Code flow — VerifySbb đóng vai trò duy nhất điều phối",
              "onEvent(VerifyRequestEvent) → resolve → verify → score → complete → CDR")
    # pipeline top row
    add_multi_box(s, Inches(0.4), Inches(1.35), Inches(1.8), Inches(1.3),
                  ["VerifyResource", "(Quarkus)"], fill=WAVE, line=OCEAN, size=11)
    add_arrow(s, Inches(2.2), Inches(2.0), Inches(2.75), Inches(2.0))
    add_multi_box(s, Inches(2.75), Inches(1.35), Inches(2.2), Inches(1.3),
                  ["submit event", "VerifyRequestEvent\ncontainer.fireEvent()"], fill=FOAM,
                  line=AQUA, size=11)
    add_arrow(s, Inches(4.95), Inches(2.0), Inches(5.5), Inches(2.0))
    add_multi_box(s, Inches(5.5), Inches(1.35), Inches(3.0), Inches(1.3),
                  ["VerifySbb.onEvent", "event mask → activity\nCMP + FSM"], fill=NAVY,
                  line=NAVY, t_color=WHITE, size=11, bold=True)
    add_arrow(s, Inches(8.5), Inches(2.0), Inches(9.05), Inches(2.0))
    add_multi_box(s, Inches(9.05), Inches(1.35), Inches(3.8), Inches(1.3),
                  ["coordinator.complete + CDR", "persist result · flow CDR\nMSISDN masked"], fill=WAVE,
                  line=OCEAN, size=11)
    # stage boxes
    stages = [
        ("1 RESOLVE", "resolverRa.sendCommand\nResolveCommand(ip, port, ts)", "PGW/PCRF/CGNAT"),
        ("2 VERIFY", "mapVerifierRa / s6aVerifierRa\nSendAuthInfo / PSI / AIR", "HLR/HSS"),
        ("3 SCORE", "AssurancePolicy\nscore ≥ threshold?", "risk class adapt"),
    ]
    xs = [0.4, 4.7, 9.0]
    for x, (nm, code, tgt) in zip(xs, stages):
        add_multi_box(s, Inches(x), Inches(3.15), Inches(4.2), Inches(1.9),
                      [nm, code, ""], fill=FOAM, line=AQUA, size=11, bold=True)
        add_textbox(s, Inches(x), Inches(4.75), Inches(4.2), Inches(0.3),
                    f"→ {tgt}", size=10, color=GRAY, align=PP_ALIGN.CENTER)
        if x < 9.0:
            end = Inches(4.7) if x < 4.7 else Inches(13.2)
            add_arrow(s, Inches(x + 4.2), Inches(4.1), end, Inches(4.1))
    add_textbox(s, Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.4),
                "Mô hình event-driven thay vì gọi trực tiếp backend:\n"
                "• REST không sở hữu trạng thái — chỉ submit event và await CompletableFuture.\n"
                "• Mỗi bước là command đi qua RA endpoint → bootstrap.fireEvent() (không bypass container).\n"
                "• Timeout dialog do timer of container cắt — abort ⇒ FALLBACK; không rò TCAP dialog.",
                size=12, color=DARK)
    add_footer(s, prs, page)
    return s


def slide_assurance(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Assurance scoring — điểm tin cậy mỗi lần verify",
              "APPROVE ⇔ score ≥ threshold VÀ resolved == claimed (khi có claimed)")
    add_multi_box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.1),
                  ["score = w₁·ipBindingFresh  +  w₂·reachable  +  w₃·notSimSwapped  +  w₄·locationPlausible"],
                  fill=NAVY, line=NAVY, t_color=WHITE, size=16, bold=True)
    weights = [
        ("w₁", "binding tươi", "IP+port+ts còn hiệu lực"),
        ("w₂", "subscriber live", "MAP PSI / Diameter REACH"),
        ("w₃", "không SIM-swap", "SAI / last-update check"),
        ("w₄", "vị trí hợp lý", "VLR/MME phù hợp tương quan"),
    ]
    y = 2.95
    for w, nm, sub in weights:
        add_multi_box(s, Inches(0.7), Inches(y), Inches(1.2), Inches(0.85),
                      [w, nm], fill=WAVE, line=OCEAN, size=12, bold=True)
        add_textbox(s, Inches(2.1), Inches(y + 0.15), Inches(2.6), Inches(0.5),
                    sub, size=11, color=GRAY)
        y += 0.95
    add_multi_box(s, Inches(5.6), Inches(2.95), Inches(7.1), Inches(3.0),
                  ["Quy tắc áp dụng", "• Ngưỡng (threshold) theo risk: login thường thấp, giao dịch giá trị cao nâng ngưỡng / step-up",
                   "• Mọi giai đoạn thiếu bằng chứng → không APPROVE (fail-closed)",
                   "• accessTech do device khai chỉ là advisory — không nâng assurance",
                   "• WIFI / FIXED (không SIM) bị từ chối ở /session-tuple — không gieo binding cellular"],
                  fill=FOAM, line=AQUA, size=11, bold=True)
    add_footer(s, prs, page)
    return s


def slide_fallback(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Fallback — quyết định theo access technology",
              "AccessTech là bearer claim; device khai WIFI/FIXED không bao giờ seed binding cellular")
    add_multi_box(s, Inches(0.4), Inches(1.5), Inches(2.2), Inches(1.0),
                  ["Auth request"], fill=NAVY, line=NAVY, t_color=WHITE, size=13, bold=True)
    add_arrow(s, Inches(2.6), Inches(2.0), Inches(3.15), Inches(2.0))
    add_multi_box(s, Inches(3.15), Inches(1.4), Inches(2.7), Inches(1.3),
                  ["AccessTech?", "GS_2G3G / LTE / NR\n(device khai + SEAM)?"], fill=FOAM,
                  line=AQUA, size=12, bold=True)
    # cellular branch
    add_arrow(s, Inches(3.15), Inches(2.75), Inches(3.15), Inches(3.6), width=1.6)
    add_arrow(s, Inches(6.0), Inches(2.05), Inches(6.55), Inches(2.05))
    add_textbox(s, Inches(5.15), Inches(1.62), Inches(1.1), Inches(0.3),
                "CÓ", size=10, bold=True, color=GRAY)
    add_multi_box(s, Inches(6.55), Inches(1.5), Inches(2.6), Inches(1.1),
                  ["Silent Auth", "IP-match: Resolver + Verifier\n(hoặc TS.43 khi có SIM)"], fill=WAVE,
                  line=OCEAN, size=11, bold=True)
    add_arrow(s, Inches(7.9), Inches(2.6), Inches(7.9), Inches(3.55), width=1.6)
    add_multi_box(s, Inches(6.55), Inches(3.6), Inches(2.6), Inches(1.5),
                  ["Scoring", "resolved==claimed?\nscore ≥ threshold?"], fill=FOAM, line=AQUA,
                  size=11, bold=True)
    add_arrow(s, Inches(6.55), Inches(4.35), Inches(5.6), Inches(4.9), width=1.6)
    add_arrow(s, Inches(7.9), Inches(5.1), Inches(7.9), Inches(6.1), width=1.6)
    add_multi_box(s, Inches(4.2), Inches(4.9), Inches(1.4), Inches(0.8),
                  ["APPROVED"], fill=RGBColor(0xE4, 0xF2, 0xE9), line=OCEAN,
                  size=11, bold=True, t_color=DARK)
    # wifi branch
    add_arrow(s, Inches(4.5), Inches(2.05), Inches(4.85), Inches(2.05))
    add_textbox(s, Inches(3.7), Inches(1.62), Inches(1.1), Inches(0.3),
                "KHÔNG", size=10, bold=True, color=GRAY)
    add_multi_box(s, Inches(10.0), Inches(1.4), Inches(2.8), Inches(1.3),
                  ["TS.43 Wi-Fi (SIM EAP-AKA)?", "entitlement + SWm/SWx"], fill=WHITE,
                  line=OCEAN, size=11, bold=True)
    add_arrow(s, Inches(12.8), Inches(2.05), Inches(12.95), Inches(2.05), head=False)
    add_multi_box(s, Inches(9.6), Inches(3.0), Inches(3.3), Inches(3.1),
                  ["FALLBACK (firewalled)", "Passkey · TOTP\nSMS OTP chỉ qua Home Routing\n"
                   "SIM-swap check trước khi gửi"], fill=RGBColor(0xFB, 0xF3, 0xE0),
                  line=SAND, size=11, bold=True, t_color=DARK)
    add_arrow(s, Inches(10.0), Inches(2.7), Inches(10.0), Inches(3.0), width=1.6)
    add_arrow(s, Inches(9.15), Inches(3.6), Inches(8.2), Inches(4.4), width=1.6, dash=True)
    add_textbox(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.5),
                "UE SDK fail-closed khi không đọc được radio; browser không pin được socket → chỉ observe. "
                "Task P-H8: cần endpoint echo để biết srcPort sau CGNAT.",
                size=11, color=GRAY)
    add_footer(s, prs, page)
    return s


def slide_verifier_msgs(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Verifier — bản đồ message với chính operator (không liên mạng)",
              "TS 29.002 (MAP) · TS 29.272 (S6a) — mỗi giai đoạn một dialog, timeout abort")
    rows = [
        ("2G/3G · MAP (intra)", "PSI", "subscriber state + location", "Cat 2.1 · được chọn"),
        ("2G/3G · MAP", "ATI", "any-time interrogation", "Cat 1 — NỘI MẠNG, cấm liên mạng"),
        ("2G/3G · MAP", "SAI", "auth vectors / SIM-swap freshness", "Cat 3.2"),
        ("LTE/5G · Diameter S6a", "AIR/AIA", "authentication info request", "FS.19"),
        ("LTE/5G · Diameter S6a", "IDR/IDA", "device/sub state (insert data)", "FS.19"),
    ]
    add_table_sea(s, ["Truy nhập", "Operation", "Mục đích", "Phân loại"], rows, Inches(0.5),
                  Inches(1.3), Inches(12.3),
                  [Inches(2.4), Inches(1.3), Inches(6.4), Inches(2.2)], row_h=0.62)
    add_textbox(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.6),
                "Các quy tắc verifier:\n"
                "• Chỉ hỏi HLR/HSS của operator nhà mình — spoofed GT không được tin (FS.11 §3.3.4).\n"
                "• jSS7 đường đi thật: AnyTimeInterrogation* / ProvideSubscriberInfo* / SendAuthenticationInfo* "
                "(MAPServiceMobility) — không ATI qua interconnect.\n"
                "• Kết quả: {reachable, notSimSwapped, locationPlausible} → đưa vào score.",
                size=12, color=DARK)
    add_footer(s, prs, page)
    return s


def slide_unified(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Kiến trúc thống nhất — Strategy A + B, bổ trợ chứ không thay thế",
              "A: thay OTP bằng silent auth · B: bảo vệ chính kênh SMS/SS7/Diameter")
    add_multi_box(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.3),
                  ["STRATEGY A — Identity / Application layer", "SAS · CAMARA /verify · TS.43 · IP-match (micro-jainslee)",
                   "Resolver + Verifier + Policy · fail-closed"], fill=NAVY,
                  line=NAVY, t_color=WHITE, size=13, bold=True)
    add_box(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.5),
            "STRATEGY B — Signalling / Interconnect border", fill=AQUA,
            t_color=WHITE, line=AQUA, size=13, bold=True)
    cols = [
        ("4G/5G", "Diameter Edge Agent\nDEA · FS.19\nSEPP/N32 · FS.36"),
        ("SS7", "SS7 Firewall\nMAP filter · FS.11\nroaming guard"),
        ("SMS", "SMS Home Routing\nA2P giữ tại operator\nMT-spoof chặn"),
        ("SAS", "Fallback OTP\nchỉ gửi khi đã\nSIM-swap check"),
    ]
    x = 0.5
    for nm, sub in cols:
        add_multi_box(s, Inches(x), Inches(3.6), Inches(2.95), Inches(1.9),
                      [nm, sub], fill=FOAM, line=OCEAN, size=12, bold=True)
        x += 3.1
    add_textbox(s, Inches(0.5), Inches(5.75), Inches(12.3), Inches(1.0),
                "Thứ tự triển khai: (1) bảo vệ SMS → (2) Diameter/5G → (3) giới thiệu silent auth → "
                "(4) dịch chuyển lưu lượng; OTP thu nhỏ về kênh fallback đã firewalled. "
                "Silent auth không thay thế SS7/Diameter firewall.",
                size=13, color=DARK)
    add_footer(s, prs, page)
    return s


def slide_security(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "An toàn — checklist không được thụt lùi",
              "Mọi điều khoản đều thuộc design ràng buộc (fail-closed, không liên mạng)")
    bullets(s, Inches(0.5), Inches(1.35), Inches(6.3), Inches(5.3), [
        "Không ATI liên mạng — SAS chỉ hỏi HLR/HSS nhà mình (FS.11 Cat 1)",
        "Fail-closed: thiếu bằng chứng không bao giờ APPROVE",
        "Idempotency: reqId dedup — một dialog MAP/Diameter mỗi giai đoạn",
        "Dialog leak: TC timer giới hạn, timeout ⇒ abort()",
        "Race: binding đọc là point-in-time (ts), không phải 'latest'",
        "Replay: bank→SAS mTLS; cửa sổ ts + reqId",
    ], size=13, gap=9)
    bullets(s, Inches(7.0), Inches(1.35), Inches(5.9), Inches(5.3), [
        "CGNAT: luôn IP+port+ts; >1 MSISDN trên cùng IP ⇒ reject",
        "Bearer: device khai WIFI/FIXED từ chối ở /session-tuple",
        "SDK fail-closed nếu không đọc được radio; không bindProcessToNetwork (rò process-wide)",
        "Privacy: MSISDN/IMSI không bao giờ trả về app — chỉ bank backend",
        "Spoofed GT: chỉ tin response từ chính HSS (FS.11 §3.3.4)",
        "CDR: masking MSISDN, IMSI không vào CDR — không đổi outcome",
    ], size=13, gap=9)
    add_footer(s, prs, page)
    return s


def slide_fs11(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "FS.11 — phân loại MAP mà silent auth tuân thủ",
              "GSMA FASG FS.11: categories của SS7 MAP — thiết kế SAS nằm trong đường được phép")
    rows = [
        ("Cat 1", "ATI quá giang", "Đọc thông tin thuê bao qua network khác", "Bị cấm — SAS tuyệt đối không dùng"),
        ("Cat 2.1", "PSI nội mạng", "Trạng thái + vị trí thuê bao", "Dùng — ProvideSubscriberInfo"),
        ("Cat 2.2", "PSI liên mạng", "Hỏi vị trí qua biên giới", "Cấm — không bao giờ có điều kiện này"),
        ("Cat 3.2", "SAI nội mạng", "Auth vectors / SIM-swap freshness", "Dùng — SendAuthenticationInfo"),
    ]
    add_table_sea(s, ["Cat", "Loại", "Ý nghĩa FS.11", "Quyết định SAS"], rows, Inches(0.5),
                  Inches(1.3), Inches(12.3),
                  [Inches(1.0), Inches(2.2), Inches(5.0), Inches(4.1)], row_h=0.75)
    add_textbox(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.4),
                "Sự phân biệt ghi trong design: P2 MAP transport thật (jSS7 coral-valley) lái PSI + SAI "
                "về đúng HLR/HSS operator — không ATI, không đường quá giang. "
                "Xem docs/research/3gpp-ts29-002-map.md + docs/design/silent-auth-flow.md.",
                size=12, color=DARK)
    add_footer(s, prs, page)
    return s


def slide_commercial(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Mô hình thương mại — Digicom-ET bán VAS, không lấy tiền SMS",
              "Lớp adapter phía trên Ethio Telecom — A2P SMS vẫn chạy qua operator")
    add_multi_box(s, Inches(0.5), Inches(1.5), Inches(3.6), Inches(1.8),
                  ["Ngân hàng", "trả tiền per /verify\nmất OTP cost, mất friction"], fill=WAVE,
                  line=OCEAN, size=13, bold=True)
    add_arrow(s, Inches(4.1), Inches(2.4), Inches(4.85), Inches(2.4))
    add_multi_box(s, Inches(4.85), Inches(1.5), Inches(3.6), Inches(1.8),
                  ["Digicom-ET SAS", "CAMARA NV2 · IP-match + TS.43\nmicro-jainslee container"], fill=NAVY,
                  line=NAVY, t_color=WHITE, size=13, bold=True)
    add_arrow(s, Inches(8.45), Inches(2.4), Inches(9.2), Inches(2.4))
    add_multi_box(s, Inches(9.2), Inches(1.5), Inches(3.6), Inches(1.8),
                  ["Ethio Telecom", "PGW · HLR/HSS · PCRF · SMSC\nkhông mất A2P SMS"], fill=FOAM,
                  line=OCEAN, size=13, bold=True)
    bullets(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(2.7), [
        "Billing qua CDR full-flow: per verify thành công + per fallback OTP (nếu dùng)",
        "Digicom thuê quyền truy nhập Resolver/Verifier từ operator — SAS không hưởng SMS revenue",
        "Bank không cần biết hạ tầng SS7 — chỉ cần mTLS + API key theo CAMARA",
        "Rủi ro trách nhiệm: silent auth là lớp phụ trợ, không thay thế firewall; phân phối theo dual license (AGPL / Operator)",
    ], size=13, gap=9)
    add_footer(s, prs, page)
    return s


def slide_roadmap(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_title(s, "Lộ trình 4 pha — SMS được bảo vệ trước, silent auth sau",
              "Chi phí thấp nhất, rủi ro thấp nhất, đưa OTP vào kênh firewalled")
    phases = [
        ("P1 · Bảo vệ SMS", "Home Routing + SS7 FW", "ngày 0–60"),
        ("P2 · Diameter/5G", "DEA FS.19 · SEPP FS.36", "60–120"),
        ("P3 · Silent Auth", "NV2/TS.43 · IP-match", "120–180"),
        ("P4 · Dịch chuyển", "OTP → fallback duy nhất", "180+"),
    ]
    x = 0.5
    for nm, sub, window in phases:
        add_multi_box(s, Inches(x), Inches(1.7), Inches(2.95), Inches(2.1),
                      [nm, sub, window], fill=FOAM, line=OCEAN, size=12, bold=True)
        x += 3.1
    add_arrow(s, Inches(3.45), Inches(2.75), Inches(3.6), Inches(2.75), head=False)
    add_arrow(s, Inches(7.0), Inches(2.75), Inches(7.5), Inches(2.75), head=False)
    add_arrow(s, Inches(10.4), Inches(2.75), Inches(10.9), Inches(2.75), head=False)
    bullets(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.7), [
        "SAS hiện tại: CAMARA NV adapter P0 xong, P2 verifier jSS7 (PSI/SAI) lab-pass — sẵn cho pilot bank",
        "Còn mở: assurance weights, post-CGNAT srcPort (P-H8), TS.43 feasibility, UAT mTLS/SS7 thật",
    ], size=13, gap=8)
    add_box(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.85),
            "Hãy bắt đầu pilot với một bank: dev sandbox + API sandbox CAMARA — 4 tuần để KPI bằng chứng.",
            fill=NAVY, t_color=WHITE, size=14, bold=True)
    add_footer(s, prs, page)
    return s


def slide_thanks(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, NAVY)
    add_wave_bottom(s, OCEAN, top=Inches(5.6))
    add_wave_bottom(s, AQUA, top=Inches(6.0))
    add_wave_bottom(s, CYAN, top=Inches(6.4))
    add_wave_bottom(s, WHITE, top=Inches(6.85))
    add_textbox(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.9),
                "Ameseginalehu · Cảm ơn", size=40, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.5),
                "Silent Authentication — đi cùng nhau để OTP không còn là khâu yếu nhất",
                size=18, color=CYAN, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.4),
                "Digicom-ET · Ethio Telecom adapter · CAMARA + GSMA FASG · nhanth87",
                size=13, color=WAVE, align=PP_ALIGN.CENTER)
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs, 1)
    slide_microjainslee(prs, 2)
    slide_h24(prs, 3)
    slide_agenda(prs, 4)
    slide_why(prs, 5)
    slide_otp_problem(prs, 6)
    slide_solution(prs, 7)
    slide_twostage(prs, 8)
    slide_e2e(prs, 9)
    slide_camara(prs, 10)
    slide_fsm(prs, 11)
    slide_code_flow(prs, 12)
    slide_assurance(prs, 13)
    slide_fallback(prs, 14)
    slide_verifier_msgs(prs, 15)
    slide_unified(prs, 16)
    slide_security(prs, 17)
    slide_fs11(prs, 18)
    slide_commercial(prs, 19)
    slide_roadmap(prs, 20)
    slide_thanks(prs, 21)

    assert len(prs.slides) == TOTAL, f"Expected {TOTAL}, got {len(prs.slides)}"
    prs.save(OUT)
    print(f"Saved → {OUT}")
    print(f"Slides: {len(prs.slides)}")
    nonempty = sum(1 for sl in prs.slides if len(sl.shapes) > 2)
    print(f"Slides with content (>2 shapes): {nonempty}/{len(prs.slides)}")


if __name__ == "__main__":
    build()