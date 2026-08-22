#!/usr/bin/env python3
"""Digicom-ET Silent Auth v3 — mix v1 story (20 slides) + v2 technical depth."""

from pathlib import Path
import subprocess
import shutil

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
ASSETS_V1 = ROOT / "assets"
ASSETS_V2 = ROOT / "assets" / "v2"
OUT = ROOT / "DigicomET_Silent_Auth_Mix_v3.pptx"
TOTAL = 28

GREEN = RGBColor(0x07, 0x89, 0x30)
YELLOW = RGBColor(0xFC, 0xDD, 0x09)
RED = RGBColor(0xDA, 0x12, 0x1A)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
CREAM = RGBColor(0xFF, 0xF8, 0xE7)
TEAL = RGBColor(0x0D, 0x73, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x4B, 0x55, 0x63)
LIGHT = RGBColor(0xE8, 0xF5, 0xE9)


def set_run(run, size=16, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=16, bold=False,
                color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_flag_bar(slide, prs):
    w = prs.slide_width
    h = Inches(0.16)
    for i, color in enumerate([GREEN, YELLOW, RED]):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(i * w // 3), 0, Emu(w // 3 + 1000), h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()


def fill_bg(slide, color=CREAM):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, prs, page):
    add_textbox(
        slide, Inches(0.4), Inches(7.15), Inches(12), Inches(0.28),
        f"Digicom-ET  ·  Silent Auth Mix v3 (story + MAP/Diameter)  ·  {page}/{TOTAL}",
        size=10, color=GRAY,
    )


def add_title(slide, title):
    add_textbox(slide, Inches(0.5), Inches(0.32), Inches(12), Inches(0.55),
                title, size=26, bold=True, color=DARK)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.88), Inches(2.4), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()


def svg_to_png(svg_path: Path, png_path: Path, scale=2):
    if png_path.exists() and png_path.stat().st_mtime >= svg_path.stat().st_mtime:
        return
    if shutil.which("rsvg-convert"):
        subprocess.run(
            ["rsvg-convert", "-z", str(scale), "-o", str(png_path), str(svg_path)],
            check=True,
        )
        return
    raise RuntimeError("rsvg-convert required")


def resolve_asset(name: str) -> Path:
    """name may be '01_foo.svg' (v1) or 'v2_bar.svg' (v2)."""
    if name.startswith("v2_"):
        svg = ASSETS_V2 / name
    else:
        svg = ASSETS_V1 / name
    if not svg.exists():
        raise FileNotFoundError(svg)
    png = svg.with_suffix(".png")
    svg_to_png(svg, png)
    return png


def add_picture(slide, name, left, top, width=None, height=None):
    png = resolve_asset(name)
    kw = {"image_file": str(png), "left": left, "top": top}
    if width is not None:
        kw["width"] = width
    if height is not None:
        kw["height"] = height
    return slide.shapes.add_picture(**kw)


def bullets(slide, left, top, width, height, items, size=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size=size, color=DARK)


def content_with_img(prs, title, items, page, img, img_w=Inches(5.5)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_flag_bar(slide, prs)
    add_title(slide, title)
    bullets(slide, Inches(0.5), Inches(1.15), Inches(6.2), Inches(5.5), items, size=14)
    add_picture(slide, img, Inches(7.0), Inches(1.2), width=img_w)
    add_footer(slide, prs, page)
    return slide


def diagram_full(prs, title, page, img, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_flag_bar(slide, prs)
    add_title(slide, title)
    add_picture(slide, img, Inches(0.45), Inches(1.1), width=Inches(12.4))
    if caption:
        add_textbox(slide, Inches(0.5), Inches(6.75), Inches(12), Inches(0.3),
                    caption, size=11, color=TEAL)
    add_footer(slide, prs, page)
    return slide


def section_slide(prs, title, subtitle, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_flag_bar(slide, prs)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.2), Inches(13.333), Inches(2.5)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = GREEN
    panel.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.9),
                title, size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.7),
                subtitle, size=16, color=CREAM, align=PP_ALIGN.CENTER)
    add_footer(slide, prs, page)
    return slide


def table_slide(prs, title, headers, rows, page, col_widths=None, note=None):
    """Full-width table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide)
    add_flag_bar(slide, prs)
    add_title(slide, title)
    n_cols = len(headers)
    n_rows = len(rows) + 1
    top = Inches(1.05)
    height = Inches(0.42) * n_rows
    if height > Inches(5.6):
        height = Inches(5.6)
    width = Inches(12.3)
    left = Inches(0.5)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        set_run(run, size=11, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        bg = LIGHT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            set_run(run, size=10, color=DARK)
    if note:
        add_textbox(slide, Inches(0.5), Inches(6.75), Inches(12), Inches(0.3),
                    note, size=11, color=TEAL)
    add_footer(slide, prs, page)
    return slide


def ensure_assets():
    # regenerate if missing
    if not (ASSETS_V1 / "01_problem_sms.svg").exists():
        subprocess.run(["python3", str(ROOT / "scripts" / "generate_svgs.py")], check=True)
    if not (ASSETS_V2 / "v2_ati_flow.svg").exists():
        subprocess.run(["python3", str(ROOT / "scripts" / "generate_svgs_v2.py")], check=True)


def build():
    ensure_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    titles = []

    # 1 Title
    titles.append("Title")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    panel = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.3), Inches(13.333), Inches(4.3)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = GREEN
    panel.line.fill.background()
    add_textbox(s, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.4),
                "DIGICOM-ET  ·  ETHIOPIA", size=18, bold=True, color=YELLOW,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.0),
                "Silent Authentication for Banks", size=38, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.5),
                "Business story + MAP/Diameter message flows (ATI · PSI · SAI · S6a)",
                size=16, color=CREAM, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.4),
                "Citizen / e-Gov login without SMS  ·  Digicom VAS adapter on Ethio Telecom",
                size=14, color=YELLOW, align=PP_ALIGN.CENTER)
    add_footer(s, prs, 1)

    # 2 Agenda
    titles.append("Agenda")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    add_title(s, "Agenda")
    bullets(s, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5), [
        "Why government & banks need Silent Auth (e-Gov, SMS OTP pain, threats)",
        "Digicom-ET offer: VAS adapter — not a competing SMSC",
        "Technical core: two-stage Resolver + Verifier (ATI / PSI / SAI / S6a)",
        "FSM, timeouts, fallback SMS, FS.11 security",
        "Standards tables: CAMARA APIs + GSMA FASG (FS.07–FS.36, SG.22)",
        "jSS7 hooks, value, Ethio Telecom partnership, roadmap",
    ], size=15)
    add_footer(s, prs, 2)

    # 3 Why Ethiopia
    titles.append("Why Ethiopia")
    content_with_img(prs, "Why Ethiopia, why now", [
        "Mobile-first banking growth across Addis and the regions",
        "SMS OTP still default — slow, fragile, costly for banks",
        "Need frictionless login without losing fraud controls",
        "Opportunity: network-side identity as VAS on Ethio Telecom",
        "Digicom-ET brings Silent Auth to government & banks",
    ], 3, "08_ethiopia.svg")

    # 4 Problem
    titles.append("SMS OTP problem")
    content_with_img(prs, "The SMS OTP problem", [
        "Customers wait 3–30s for a code — many abandon login",
        "SS7 / interconnect can intercept or redirect MT-SMS",
        "SIM-swap: attacker receives OTP on a new SIM",
        "Phishing / real-time OTP proxies steal typed codes",
        "Banks pay SMS costs and still absorb fraud losses",
    ], 4, "01_problem_sms.svg")

    # 5 Persona
    titles.append("Government digital services")
    content_with_img(prs, "Government digital services", [
        "Citizens access e-Gov portals: tax, civil registry, social payments",
        "Want: open app → prove phone number → done",
        "Do not want: wait for SMS OTP, type 6 digits, fail delivery",
        "Must stay safe against SIM-swap and SS7 intercept",
        "Digicom Silent Auth for government + bank partners",
    ], 5, "02_persona.svg")

    # 6 Threats
    titles.append("Threats")
    content_with_img(prs, "Threats Silent Auth removes", [
        "SS7 intercept — no SMS payload to steal",
        "Phishing / OTP relay — nothing for the user to leak",
        "Premium-rate / AIT abuse — no OTP blast to fraud numbers",
        "Delivery failure — login not hostage to SMS latency",
        "SIM-swap — SAI / lastUpdateLocation freshness check",
    ], 6, "06_threats.svg")

    # 7 Solution divider
    titles.append("The solution")
    section_slide(
        prs, "The solution",
        "Digicom-ET Silent Auth = business VAS + MAP/Diameter verification",
        7,
    )

    # 8 Two-stage (v2 tech into v1 story)
    titles.append("Two-stage design")
    content_with_img(prs, "Hard fact: two stages (not one MAP call)", [
        "MAP / Diameter cannot map IP → MSISDN",
        "Stage 1 RESOLVER: IP:port:ts → MSISDN via PGW/PCRF/CGNAT",
        "Stage 2 VERIFIER: ATI / PSI / SAI or Diameter S6a → assurance",
        "CGNAT ⇒ app MUST send IP + source port + timestamp",
        "Reject if resolver returns >1 MSISDN → FALLBACK",
        "Happy path: Approve in ≤3s — no SMS, no OTP typing",
    ], 8, "v2_two_stage.svg", img_w=Inches(5.8))

    # 9 E2E flow
    titles.append("E2E message flow")
    diagram_full(
        prs, "End-to-end message flow", 9, "v2_e2e_sequence.svg",
        caption="App → Bank BE → Digicom SAS → PGW resolve → MAP/Diameter verify → HLR/HSS → Approve",
    )

    # 10 ATI deep dive
    titles.append("ATI deep dive")
    content_with_img(prs, "ATI — AnyTimeInterrogation (MAP)", [
        "Verifier opens TC dialog → sends ATI to home HLR",
        "Response: subscriberState, locationInfo, VLR address",
        "FS.11 Category 1 on interconnect → BLOCKED externally",
        "Deployment invariant: Digicom SAS inside operator only",
        "Queries own HLR/HSS — never cross-operator ATI",
        "jSS7: AnyTimeInterrogationRequestImpl (coral-valley)",
    ], 10, "v2_ati_flow.svg", img_w=Inches(5.8))

    # 11 PSI / SAI / Diameter
    titles.append("PSI · SAI · Diameter")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    add_title(s, "Verifier messages: PSI · SAI · Diameter S6a")
    bullets(s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.5), [
        "PSI (ProvideSubscriberInfo) — Cat 2.1: state + location",
        "SAI (SendAuthenticationInfo) — Cat 3.2: SIM-swap freshness",
        "Prefer PSI for reachability; SAI when swap risk is high",
        "4G/5G: Diameter S6a IDR/IDA + AIR/AIA (FS.19)",
        "One dialog max per request; abort on timeout",
        "jSS7: ProvideSubscriberInfoRequestImpl, SendAuthenticationInfoRequestImpl",
    ], size=13)
    add_picture(s, "v2_psi_sai.svg", Inches(6.7), Inches(1.15), width=Inches(6.1))
    add_footer(s, prs, 11)

    # 12 Architecture + money
    titles.append("Architecture & pocket")
    content_with_img(prs, "Adapter on top — we do NOT take telco SMS revenue", [
        "Digicom-ET = thin VAS between bank backends and Ethio Telecom",
        "Reuses PGW, HLR/HSS, SMSC — no core rewrite",
        "Banks pay Digicom per /verify — not SMS wholesale",
        "A2P / interconnect SMS stays 100% with the operator",
        "Fallback OTP still billed through Ethio Telecom SMSC",
        "Optional API revenue share — additive, not cannibalizing",
    ], 12, "04_architecture.svg", img_w=Inches(5.6))

    # 13 Fallback
    titles.append("Fallback SMS")
    content_with_img(prs, "Fallback: Silent first → SMS OTP if needed", [
        "Wi-Fi-only, no binding, timeout, low assurance → FALLBACK",
        "Fail-closed: never soft-approve without evidence",
        "SMS is safety net — silent path is primary",
        "Digicom orchestrates policy; SMS rides operator rails",
        "Residual OTP can sit behind SMS Home Routing / SS7 FW",
    ], 13, "05_fallback.svg", img_w=Inches(5.4))

    # 14 FSM + timeout
    titles.append("FSM & timeouts")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    add_title(s, "SAS FSM + timeout / dialog-leak strategy")
    bullets(s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.5), [
        "RESOLVING → VERIFYING → SCORING → APPROVED",
        "Any missing evidence → FALLBACK (fail-closed)",
        "Resolver budget: 300 ms",
        "MAP / Diameter dialog: 2 s — then abort()",
        "Total SAS budget: ≤ 3 s",
        "SAS is dialog anchor — no hung HSS stalls the bank app",
    ], size=13)
    add_picture(s, "v2_fsm.svg", Inches(6.7), Inches(1.1), width=Inches(6.0))
    add_footer(s, prs, 14)

    # 15 Security + defense
    titles.append("Security")
    content_with_img(prs, "Security checklist + defense in depth", [
        "No interconnect ATI (FS.11 Cat.1) — intra-HLR only",
        "Idempotent reqId; mTLS bank→SAS; ts window anti-replay",
        "MSISDN/IMSI never returned to the mobile app",
        "Layer 1 Silent Auth · Layer 2 SMS FW · Layer 3 SS7/Diameter/SEPP",
        "Aligned with GSMA FS.11 / FS.19 / FS.21 / FS.36",
    ], 15, "10_shield.svg", img_w=Inches(5.2))

    # 16 Integration + jSS7
    titles.append("Integration & jSS7")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    add_title(s, "Bank integration + jSS7 hooks (coral-valley)")
    left = [
        "1. Enable Digicom SDK in bank app",
        "2. App sends deviceCred + IP:port:ts to bank BE",
        "3. BE → Digicom POST /verify",
        "4. Digicom returns match + assurance (or FALLBACK)",
        "5. Bank issues session — or triggers SMS OTP",
    ]
    right = [
        "AnyTimeInterrogationRequestImpl — ATI",
        "ProvideSubscriberInfoRequestImpl — PSI",
        "SendAuthenticationInfoRequestImpl — SAI",
        "MAPProviderImpl → MAPServiceMobility",
        "Open: jDiameter S6a client module",
    ]
    bullets(s, Inches(0.5), Inches(1.15), Inches(6.0), Inches(5.5), left, size=13)
    box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.2), Inches(5.9), Inches(5.2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = TEAL
    add_textbox(s, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.4),
                "jSS7 MAP classes", size=14, bold=True, color=TEAL)
    bullets(s, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.2), right, size=12)
    add_footer(s, prs, 16)

    # ---- Standards block (slides 17–23) ----
    titles.append("Standards section")
    section_slide(
        prs, "Standards & APIs",
        "CAMARA Open Gateway  ·  GSMA FASG security (FS / SG / FF)",
        17,
    )

    # 18 CAMARA table
    titles.append("CAMARA APIs")
    table_slide(
        prs, "CAMARA / Open Gateway APIs (identity & fraud)",
        ["API", "What it does", "Digicom role"],
        [
            ("Number Verification (NV)", "Match claimed MSISDN to live cellular session", "Core Silent Auth /verify"),
            ("Number Verification 2 / TS.43", "SIM EAP-AKA — works on Wi-Fi too", "Phase 3 coverage expand"),
            ("SIM Swap", "Detect recent SIM / number change", "Assurance / step-up signal"),
            ("OTP SMS", "Legacy OTP delivery (fallback path)", "Orchestrate only — SMSC = telco"),
            ("Scam Signal", "Fraud / scam indicators on number", "Risk policy input"),
            ("KYC Match", "Match subscriber KYC attributes", "e-Gov / bank onboarding"),
            ("Number Recycling", "Detect recycled MSISDN", "Avoid wrong-person auth"),
        ],
        18,
        col_widths=[Inches(3.2), Inches(5.5), Inches(3.6)],
        note="Digicom implements / adapts these for Ethio Telecom + government & bank apps",
    )

    # 19 CAMARA mapping to SAS
    titles.append("CAMARA ↔ Digicom SAS")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    add_title(s, "How CAMARA maps onto Digicom SAS")
    bullets(s, Inches(0.5), Inches(1.15), Inches(12), Inches(5.5), [
        "NV /verify  ≡  Digicom two-stage: PGW resolve + MAP/Diameter verify",
        "SIM Swap API  ≡  SAI / lastUpdateLocation age + HSS signals",
        "OTP SMS  ≡  FALLBACK branch — Digicom policy, Ethio Telecom SMSC bills",
        "TS.43 NV2  ≡  future Wi-Fi silent path (SIM credential, not bearer IP)",
        "KYC Match / Number Recycling  ≡  optional e-Gov onboarding checks",
        "App-facing contract: CAMARA-aligned HTTPS; signalling stays operator-internal",
    ], size=14)
    add_footer(s, prs, 19)

    # 20 GSMA FASG index
    titles.append("GSMA FASG index")
    table_slide(
        prs, "GSMA FASG — signalling & SMS security documents",
        ["PRD", "Title", "Scope", "Digicom relevance"],
        [
            ("FS.07", "SS7 & SIGTRAN Security", "Threats / attack methods", "SS7 SMS threat model"),
            ("FS.11", "SS7 Interconnect FW", "MAP/CAMEL Cat 1/2/3 rules", "ATI Cat.1 · SRI-SM · MT-spoof"),
            ("FS.19", "Diameter Interconnect", "LTE/5G Diameter attacks", "S6a/S6c/SGd SMS path"),
            ("FS.20", "GTP Security", "GRX/IPX GTP-C/U", "Bearer / PGW context"),
            ("FS.21", "Interconnect Recommendations", "Umbrella categorise→filter", "Policy consistency"),
            ("FS.31", "Baseline Security Controls", "Roaming / interconnect catalogue", "Border checklist"),
            ("FS.36", "5G Interconnect Security", "SEPP / N32 / PRINS", "5G SMSF + N32"),
            ("SG.22", "SMS Firewall Best Practices", "SMS content / AIT policy", "Fallback OTP hygiene"),
            ("FF.09", "SMS Fraud", "Fraud taxonomy", "SIM-swap / spam context"),
        ],
        20,
        col_widths=[Inches(1.2), Inches(3.4), Inches(3.6), Inches(4.1)],
        note="Most members-only; FS.11 v4.0 publicly circulated. Digicom aligns SAS + residual SMS protection.",
    )

    # 21 FS.11 categories
    titles.append("FS.11 MAP categories")
    table_slide(
        prs, "FS.11 — MAP packet categories (Silent Auth impact)",
        ["Cat", "Meaning", "Examples", "Digicom rule"],
        [
            ("1", "Unauthorised on interconnect — BLOCK", "ATI, SendIMSI, unknown opcode", "ATI only to OWN HLR (intra-net)"),
            ("2.1", "Operator traffic — needs answer", "PSI, PRN, PSL", "Primary verifier (state+location)"),
            ("2.2", "Operator traffic — no answer", "ISD, DSD", "Filter IMSI↔SCCP"),
            ("3.1", "Inter-op + VLR/SGSN check", "MO-FSM, USSD, IDP", "Home Routing / SMS FW"),
            ("3.2", "Inter-op + time/location", "UL, SAI", "SIM-swap freshness (SAI)"),
            ("3.3", "IPSM-GW / SMS-specific", "SMS checks", "Fallback OTP path"),
        ],
        21,
        col_widths=[Inches(1.0), Inches(4.0), Inches(3.8), Inches(3.5)],
        note="Silent Auth Verifier lives inside operator — never fires Cat.1 ATI over interconnect",
    )

    # 22 Strategy A vs B + GSMA
    titles.append("Two strategies + standards")
    table_slide(
        prs, "Replace OTP vs Protect OTP — which standard applies",
        ["Strategy", "Mechanism", "Primary standards", "Owner"],
        [
            ("A — Replace OTP", "Silent Auth (NV / TS.43)", "CAMARA NV, TS.43, FS.11/19 (intra)", "Digicom SAS"),
            ("B — Protect OTP", "Home Routing + signalling FW", "FS.11, FS.19, FS.36, SG.22, FF.09", "Telco + Digicom policy"),
            ("Umbrella", "Categorise → monitor → filter", "FS.21 + FS.31 baseline", "Operator border"),
            ("5G path", "SEPP / N32 / SMSF", "FS.36, 3GPP TS 33.501", "Operator + Digicom"),
        ],
        22,
        col_widths=[Inches(2.4), Inches(3.5), Inches(4.0), Inches(2.4)],
        note="Every OTP Digicom cannot eliminate must still be protected (Strategy B)",
    )

    # 23 Open Gateway positioning
    titles.append("Open Gateway")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    add_title(s, "Ethiopia Open Gateway positioning")
    bullets(s, Inches(0.5), Inches(1.15), Inches(12), Inches(5.5), [
        "Digicom exposes CAMARA-shaped APIs to government & banks",
        "Ethio Telecom remains network source of truth (HLR/HSS/PGW/SMSC)",
        "Compliance story: FS.11/19/21/36 for signalling; CAMARA for app contract",
        "No SMS revenue taken — Open Gateway VAS layer on top of the operator",
        "Roadmap: NV pilot → SIM Swap signal → TS.43 / NV2 Wi-Fi path",
    ], size=14)
    add_footer(s, prs, 23)

    # 24 Value government & banks
    titles.append("Value for government & banks")
    content_with_img(prs, "Value for government & banks", [
        "Higher login conversion on e-Gov and banking apps",
        "Lower SMS OTP spend on the silent happy path",
        "Stronger posture vs SS7 intercept & phishing",
        "CAMARA-aligned API — portable across Open Gateway markets",
        "Risk-based step-up for high-value / benefit transfers",
    ], 24, "09_before_after.svg", img_w=Inches(5.4))

    # 25 Telecom + roadmap
    titles.append("Telecom & roadmap")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    add_title(s, "Ethio Telecom partnership + roadmap")
    bullets(s, Inches(0.5), Inches(1.15), Inches(6.0), Inches(5.5), [
        "New enterprise VAS without cannibalizing SMS wholesale",
        "Government & banks stay on cellular — Digicom prefers data bearer",
        "Digicom sells/integrates; telco supplies network truth",
        "Phase 1: pilot IP-match + MAP ATI/PSI",
        "Phase 2: Diameter S6a + CAMARA SIM Swap",
        "Phase 3: TS.43 / NV2 Wi-Fi · FS.36 5G path",
    ], size=13)
    add_picture(s, "07_money_model.svg", Inches(6.7), Inches(1.2), width=Inches(6.0))
    add_footer(s, prs, 25)

    # 26 Open items
    titles.append("Open items")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    add_title(s, "Open items")
    bullets(s, Inches(0.5), Inches(1.15), Inches(12), Inches(5.5), [
        "Resolver source: PGW RADIUS vs PCRF Sd vs CGNAT log",
        "jDiameter S6a client (mirror jSS7 MAP verifier)",
        "CAMARA Number Verification adapter: SAS /verify ↔ NV contract",
        "Assurance weights + per-risk thresholds (e-Gov vs bank transfer)",
        "Shared identity-policy store between SAS and signalling FW",
    ], size=14)
    add_footer(s, prs, 26)

    # 27 CTA
    titles.append("CTA")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    panel = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.6)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = TEAL
    panel.line.fill.background()
    add_textbox(s, Inches(1.8), Inches(2.2), Inches(9.7), Inches(0.8),
                "Let's make government login silent", size=30, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.8), Inches(3.2), Inches(9.7), Inches(1.2),
                "Pilot Digicom-ET Silent Auth for e-Gov and banks.\n"
                "CAMARA-aligned · MAP/Diameter verify · GSMA FS.11/19 · no telco pocket touched.",
                size=15, color=CREAM, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.8), Inches(4.6), Inches(9.7), Inches(0.4),
                "contact@digicom-et.example  ·  Addis Ababa",
                size=13, color=YELLOW, align=PP_ALIGN.CENTER)
    add_footer(s, prs, 27)

    # 28 Thank you
    titles.append("Thank you")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_flag_bar(s, prs)
    add_textbox(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.0),
                "Ameseginalehu  ·  Thank you", size=38, bold=True,
                color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.5),
                "Digicom-ET — Silent Authentication for Government & Banks",
                size=16, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.4),
                "CAMARA + GSMA FASG  ·  ATI/PSI/SAI for the engineers",
                size=13, color=TEAL, align=PP_ALIGN.CENTER)
    add_footer(s, prs, 28)

    assert len(prs.slides) == TOTAL, f"Expected {TOTAL}, got {len(prs.slides)}"
    prs.save(OUT)
    print(f"Saved → {OUT}")
    print(f"Slides: {len(prs.slides)}")
    for i, t in enumerate(titles, 1):
        print(f"  {i:2d}. {t}")


if __name__ == "__main__":
    build()
