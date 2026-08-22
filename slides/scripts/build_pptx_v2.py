#!/usr/bin/env python3
"""Build Digicom-ET Silent Auth TECHNICAL v2 PowerPoint (~17 slides)."""

from pathlib import Path
import subprocess
import shutil
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets" / "v2"
OUT = ROOT / "DigicomET_Silent_Auth_Technical_v2.pptx"
TOTAL = 17

# Ethiopia palette
GREEN = RGBColor(0x07, 0x89, 0x30)
YELLOW = RGBColor(0xFC, 0xDD, 0x09)
RED = RGBColor(0xDA, 0x12, 0x1A)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
CREAM = RGBColor(0xFF, 0xF8, 0xE7)
TEAL = RGBColor(0x0D, 0x73, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x4B, 0x55, 0x63)
LIGHT = RGBColor(0xE8, 0xF5, 0xE9)

EXPECTED_SVGS = [
    "v2_two_stage.svg",
    "v2_e2e_sequence.svg",
    "v2_ati_flow.svg",
    "v2_psi_sai.svg",
    "v2_diameter_s6a.svg",
    "v2_fsm.svg",
    "v2_timeout.svg",
    "v2_adapter.svg",
]

SVG_GENERATORS = [
    "generate_svgs_v2.py",
    "generate_v2_svgs.py",  # fallback placeholder generator
]


def set_run(run, size=16, bold=False, color=DARK, font="Calibri", italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=16, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_flag_bar(slide, prs):
    w = prs.slide_width
    h = Inches(0.14)
    for i, color in enumerate([GREEN, YELLOW, RED]):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(i * w // 3), 0, Emu(w // 3 + 1000), h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()


def fill_slide_bg(slide, color=CREAM):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, prs, page):
    add_textbox(
        slide, Inches(0.4), Inches(7.15), Inches(10), Inches(0.3),
        f"Digicom-ET  ·  Silent Auth Technical Design  ·  {page}/{TOTAL}",
        size=9, color=GRAY, align=PP_ALIGN.LEFT,
    )


def add_title_bar(slide, title):
    add_textbox(slide, Inches(0.5), Inches(0.28), Inches(12), Inches(0.55),
                title, size=26, bold=True, color=DARK)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.82), Inches(2.2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()


def svg_to_png(svg_path: Path, png_path: Path, scale=2):
    if png_path.exists() and png_path.stat().st_mtime >= svg_path.stat().st_mtime:
        return png_path

    try:
        import cairosvg
        cairosvg.svg2png(url=str(svg_path), write_to=str(png_path), scale=scale)
        return png_path
    except Exception:
        pass

    if shutil.which("rsvg-convert"):
        subprocess.run(
            ["rsvg-convert", "-z", str(scale), "-o", str(png_path), str(svg_path)],
            check=True,
        )
        return png_path

    if shutil.which("inkscape"):
        subprocess.run(
            ["inkscape", str(svg_path), f"--export-filename={png_path}",
             f"--export-dpi={96 * scale}"],
            check=True,
        )
        return png_path

    try:
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM
        drawing = svg2rlg(str(svg_path))
        renderPM.drawToFile(drawing, str(png_path), fmt="PNG")
        return png_path
    except Exception as e:
        raise RuntimeError(
            f"Cannot rasterize {svg_path.name}. Install rsvg-convert or cairosvg. ({e})"
        )


def add_picture(slide, name, left, top, width=None, height=None):
    svg = ASSETS / name
    if not svg.exists():
        raise FileNotFoundError(f"Missing SVG: {svg}")
    png = ASSETS / (svg.stem + ".png")
    svg_to_png(svg, png)
    kwargs = {"image_file": str(png), "left": left, "top": top}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(**kwargs)


def _run_svg_gen():
    for name in SVG_GENERATORS:
        gen = ROOT / "scripts" / name
        if gen.exists():
            subprocess.run(["python3", str(gen)], check=True)
            return True
    return False


def ensure_svgs(retry_once=True):
    ASSETS.mkdir(parents=True, exist_ok=True)
    missing = [n for n in EXPECTED_SVGS if not (ASSETS / n).exists()]
    if missing:
        _run_svg_gen()
        missing = [n for n in EXPECTED_SVGS if not (ASSETS / n).exists()]
    if missing and retry_once:
        time.sleep(1)
        _run_svg_gen()
        missing = [n for n in EXPECTED_SVGS if not (ASSETS / n).exists()]
    if missing:
        raise FileNotFoundError(f"Missing v2 SVG assets: {missing}")


def add_bullet_paragraph(tf, text, first=False, size=15, space_after=8):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(space_after)
    p.level = 0
    run = p.add_run()
    run.text = f"•  {text}"
    set_run(run, size=size, color=DARK)
    return p


def content_slide(prs, title, bullets, page, img=None, img_left=Inches(6.8),
                  img_w=Inches(6.0), img_top=Inches(1.15), text_w=Inches(6.0)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    add_flag_bar(slide, prs)
    add_title_bar(slide, title)

    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), text_w, Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        add_bullet_paragraph(tf, b, first=(i == 0), size=15)

    if img:
        add_picture(slide, img, img_left, img_top, width=img_w)
    add_footer(slide, prs, page)
    return slide


def diagram_slide(prs, title, page, img, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    add_flag_bar(slide, prs)
    add_title_bar(slide, title)
    add_picture(slide, img, Inches(0.45), Inches(1.05), width=Inches(12.4))
    if caption:
        add_textbox(slide, Inches(0.5), Inches(6.55), Inches(12), Inches(0.4),
                    caption, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide, prs, page)
    return slide


def table_slide(prs, title, headers, rows, page, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    add_flag_bar(slide, prs)
    add_title_bar(slide, title)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    left, top = Inches(0.5), Inches(1.1)
    width = Inches(12.3)
    height = Inches(0.45 * n_rows)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        set_run(run, size=12, bold=True, color=WHITE)

    for i, row in enumerate(rows, start=1):
        bg = LIGHT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            set_run(run, size=11, color=DARK)

    add_footer(slide, prs, page)
    return slide


def build():
    ensure_svgs()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    titles = []

    # 1 — Title
    titles.append("Digicom-ET Silent Authentication — Technical Design (Ethiopia)")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(s)
    add_flag_bar(s, prs)
    panel = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.4), Inches(13.333), Inches(3.8)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = GREEN
    panel.line.fill.background()
    add_textbox(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.4),
                "DIGICOM-ET", size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(2.2), Inches(12), Inches(1.0),
                "Silent Authentication — Technical Design", size=34, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(3.3), Inches(12), Inches(0.5),
                "Two-stage resolver + MAP/Diameter verifier  ·  Ethiopia deployment",
                size=16, color=CREAM, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(4.1), Inches(12), Inches(0.35),
                "Protocol flows  ·  Fail-closed FSM  ·  GSMA FS.11 / FS.19",
                size=13, color=YELLOW, align=PP_ALIGN.CENTER)
    add_footer(s, prs, 1)

    # 2 — Scope / non-goals
    titles.append("Scope & Non-Goals")
    content_slide(prs, "Scope & Non-Goals", [
        "IN SCOPE: auth without SMS OTP on cellular data bearer",
        "IN SCOPE: VAS adapter deployed on / inside operator network",
        "IN SCOPE: bank backend calls SAS POST /verify server-to-server",
        "NOT replacing SS7 / Diameter signalling firewall (Strategy B)",
        "NOT taking operator SMS revenue — fallback OTP stays operator-billed",
        "NOT cross-operator ATI — intra-network HLR/HSS only (FS.11 Cat.1)",
    ], 2)

    # 3 — Hard constraint
    titles.append("Hard Constraint: Two Stages Required")
    content_slide(prs, "Hard Constraint: Two Stages Required", [
        "MAP / Diameter cannot map IP → MSISDN — protocol limitation",
        "Stage 1 RESOLVER: PGW/GGSN session binding → MSISDN/IMSI",
        "Stage 2 VERIFIER: MAP ATI/PSI/SAI or Diameter S6a IDR/AIR",
        "CGNAT ⇒ app MUST send IP + port + timestamp",
        "Reject if resolver returns >1 MSISDN for same binding",
    ], 3, img="v2_two_stage.svg", img_left=Inches(6.6), img_w=Inches(6.2),
       text_w=Inches(5.8))

    # 4 — Actors
    titles.append("Actors")
    table_slide(
        prs, "Actors",
        ["Actor", "Role"],
        [
            ("Bank App (mobile)", "Cellular data; collects {srcIP, srcPort, ts, deviceCred}"),
            ("Bank Backend", "Login decision; calls SAS POST /verify"),
            ("Digicom SAS", "Resolver + Verifier + Policy engine"),
            ("PGW Resolver", "PGW/GGSN / PCRF / CGNAT log → MSISDN+IMSI"),
            ("MAP/Diameter Verifier", "jSS7 (2G/3G) + jDiameter S6a (4G/5G)"),
            ("HLR / HSS", "Operator subscriber DB — intra-network only"),
        ],
        4,
        col_widths=[Inches(2.8), Inches(9.5)],
    )

    # 5 — E2E flow
    titles.append("End-to-End Message Flow")
    diagram_slide(
        prs, "End-to-End Message Flow", 5, "v2_e2e_sequence.svg",
        caption="Numbered steps match silent-auth-flow.md — happy path, cellular bearer",
    )

    # 6 — Stage 1 Resolver
    titles.append("Stage 1: IP Resolver")
    content_slide(prs, "Stage 1: IP Resolver", [
        "Input: {srcIP, srcPort, ts} from SAS",
        "Sources: PGW/GGSN Gi/SGi accounting, PCRF Gx/Sd, CGNAT log",
        "Output: MSISDN + IMSI + bearerAge (point-in-time at ts)",
        "Reject if >1 MSISDN — CGNAT ambiguity → FALLBACK",
        "No binding / Wi-Fi-only / stale bearer → FALLBACK",
        "Budget: 300 ms — on expiry → FALLBACK",
    ], 6)

    # 7 — Stage 2 Verifier
    titles.append("Stage 2: MAP/Diameter Verifier")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)
    add_flag_bar(slide, prs)
    add_title_bar(slide, "Stage 2: MAP/Diameter Verifier")
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.35),
                "One MAP/Diameter dialog max per request. SIM-swap = lastUpdateLocation age vs request time.",
                size=12, color=TEAL)
    table = slide.shapes.add_table(6, 4, Inches(0.5), Inches(1.45), Inches(12.3), Inches(2.8)).table
    headers = ["Access", "Message", "Purpose", "FS Category"]
    rows = [
        ("2G/3G", "PSI", "subscriber state + location, intra-net", "Cat 2.1"),
        ("2G/3G", "ATI", "any-time interrogation (intra-net ONLY)", "Cat 1 interconnect"),
        ("2G/3G", "SAI", "auth vectors / SIM-swap freshness", "Cat 3.2"),
        ("4G/5G", "IDR/IDA", "insert/inspect subscriber data", "FS.19"),
        ("4G/5G", "AIR/AIA", "authentication info request/answer", "FS.19"),
    ]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = h
        set_run(run, size=11, bold=True, color=WHITE)
    for i, row in enumerate(rows, start=1):
        bg = LIGHT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = val
            set_run(run, size=10, color=DARK, bold=(j == 1))
    add_footer(slide, prs, 7)

    # 8 — ATI deep dive
    titles.append("ATI Deep Dive (FS.11 Cat.1)")
    diagram_slide(
        prs, "ATI Deep Dive (FS.11 Cat.1)", 8, "v2_ati_flow.svg",
        caption="TC dialog · ATI request/response fields · deploy Digicom inside operator",
    )

    # 9 — PSI vs SAI
    titles.append("PSI vs SAI")
    diagram_slide(prs, "PSI vs SAI", 9, "v2_psi_sai.svg")

    # 10 — Diameter S6a
    titles.append("Diameter S6a Path")
    diagram_slide(prs, "Diameter S6a Path (4G/5G)", 10, "v2_diameter_s6a.svg",
                  caption="IDR/IDA + AIR/AIA — GSMA FS.19")

    # 11 — SAS FSM
    titles.append("SAS Request FSM")
    diagram_slide(prs, "SAS Request FSM — Fail-Closed", 11, "v2_fsm.svg",
                  caption="No partial approvals — any missing evidence → FALLBACK")

    # 12 — Timeout
    titles.append("Timeout & Dialog Leak Strategy")
    content_slide(
        prs, "Timeout & Dialog Leak Strategy", [
            "Resolver: 300 ms → FALLBACK",
            "MAP dialog (PSI/ATI): 2 s TC timer → abort() → FALLBACK",
            "Diameter S6a (IDR/AIR): 2 s → FALLBACK",
            "Total SAS budget: 3 s — bank shows normal login",
            "SAS is dialog anchor — never let hung HSS query stall the app",
            "Every MAP dialog has bounded TC timer — no dialog leak",
        ],
        12, img="v2_timeout.svg", img_left=Inches(6.5), img_w=Inches(6.3),
        text_w=Inches(5.8),
    )

    # 13 — Fallback SMS
    titles.append("Fallback SMS Policy")
    content_slide(prs, "Fallback SMS Policy", [
        "Trigger: Wi-Fi-only, resolver miss, MAP/Diameter timeout, low assurance",
        "SAS returns FALLBACK — bank triggers step-up (TOTP / Passkey / SIM-OTP)",
        "SMS OTP is safety net — silent path is primary",
        "Operator bills every fallback SMS via own SMSC",
        "Digicom orchestrates policy only — SMS rides operator rails",
        "Fail-closed: never approve without cryptographic network evidence",
    ], 13)

    # 14 — Security checklist
    titles.append("Security Checklist")
    content_slide(prs, "Security Checklist", [
        "No interconnect ATI — FS.11 Cat.1 blocked; intra-network HLR only",
        "Fail-closed — missing evidence never approves",
        "Idempotent reqId — dedup retries; one dialog per stage",
        "MSISDN/IMSI never returned to mobile app — bank backend only",
        "SIM-swap check via SAI / lastUpdateLocation age",
        "CGNAT: IP+port+ts required; reject >1 MSISDN; point-in-time read at ts",
        "Replay: mTLS bank→SAS; reqId + ts window; spoofed GT rejected",
    ], 14)

    # 15 — jSS7 hooks
    titles.append("jSS7 Hooks (coral-valley)")
    content_slide(prs, "jSS7 Hooks (coral-valley)", [
        "AnyTimeInterrogationRequestImpl — ATI (subscriberInformation)",
        "ProvideSubscriberInfoRequestImpl — PSI (subscriberInformation)",
        "SendAuthenticationInfoRequestImpl — SAI (authentication)",
        "MAPProviderImpl → MAPServiceMobility — dialog open / TC timer",
        "Dialog lifecycle: open → send → await → close or abort on timeout",
        "jDiameter S6a client module — open item (mirror MAP verifier)",
    ], 15)

    # 16 — Summary
    titles.append("Summary")
    content_slide(prs, "Summary — One Page", [
        "Two stages: IP:port:ts → RESOLVER → MSISDN → VERIFIER → assurance",
        "ATI intra-network only (FS.11 Cat.1) — Digicom SAS inside operator",
        "Fail-closed FSM → FALLBACK SMS (operator-billed) / step-up MFA",
        "Adapter: Digicom VAS on Ethio Telecom HLR/HSS/PGW — no SMS revenue taken",
        "Budgets: resolver 300 ms · MAP/Diameter 2 s · total ≤ 3 s",
        "jSS7: ATI/PSI/SAI; open: jDiameter S6a + resolver source",
    ], 16, img="v2_adapter.svg", img_left=Inches(6.6), img_w=Inches(6.2),
       text_w=Inches(5.8))

    # 17 — Open items
    titles.append("Open Items")
    content_slide(prs, "Open Items", [
        "Resolver interface: PGW RADIUS accounting vs PCRF Sd vs CGNAT log",
        "jDiameter S6a client module (mirror jSS7 MAP verifier)",
        "CAMARA Number Verification adapter: SAS /verify → CAMARA contract",
        "Assurance weights + per-risk thresholds (login vs money transfer)",
        "Shared identity-policy store between SAS and signalling firewall",
    ], 17)

    assert len(titles) == TOTAL, f"Expected {TOTAL} slides, got {len(titles)}"
    assert len(prs.slides) == TOTAL, f"Slide count mismatch: {len(prs.slides)}"
    prs.save(OUT)
    print(f"Saved → {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Titles:")
    for i, t in enumerate(titles, 1):
        print(f"  {i:2d}. {t}")


if __name__ == "__main__":
    build()
